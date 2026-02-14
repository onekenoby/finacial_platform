import os
import json
import uuid
import glob
import shutil
import time
import datetime
import hashlib
import base64
import re
import psycopg2
from psycopg2.extras import Json
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from typing import List, Dict, Tuple
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient, models
from neo4j import GraphDatabase
from openai import OpenAI

# --- CONFIGURAZIONE ---
BASE_DATA_DIR = "./data_ingestion"
PROCESSED_DIR = "./data_ingestion/processed"

# DB Configs
PG_DSN = "dbname=ai_ingestion user=admin password=admin_password host=localhost"

QDRANT_HOST = "localhost"
QDRANT_COLLECTION = "financial_docs"
NEO4J_URI = "bolt://localhost:7687"
NEO4J_AUTH = ("neo4j", "password_sicura")

# AI Configs
LM_STUDIO_URL = "http://localhost:1234/v1"
LM_STUDIO_API_KEY = "lm-studio"
MODEL_LLM_ID = "gemma-3-12b"
MODEL_VISION_ID = "qwen3-vl-8b"
MODEL_EMBEDDING_ID = "BAAI/bge-m3"

# --- DEFINIZIONE DINAMICA DELLE ONTOLOGIE ---

ONTOLOGIES = {
    "financial": {
        "role": "Senior Data Engineer specializing in Finance",
        "nodes": '"Organization", "Person", "Metric", "Product", "Date", "Location"',
        "edges": '"CEO_OF", "COMPETITOR_OF", "SUPPLIER_OF", "REPORTED_VALUE", "LAUNCHED", "OPERATES_IN"',
        "special_rule": "Translate financial metrics to English (e.g., 'Fatturato' -> 'Revenue')."
    },
    "educational": {
        "role": "Economic Theory Expert",
        "nodes": '"Concept", "AssetClass", "Indicator", "EconomicEvent"',
        "edges": '"CAUSES", "POSITIVELY_CORRELATED_WITH", "NEGATIVELY_CORRELATED_WITH", "HEDGES_AGAINST", "DEFINED_AS"',
        "special_rule": "Focus on causal relationships and definitions. Translate concepts to standard English terminology."
    },
    "client": {
        "role": "Private Banker / Wealth Manager",
        "nodes": '"Client", "Goal", "RiskProfile", "Asset", "FamilyMember"',
        "edges": '"HAS_GOAL", "HAS_RISK_PROFILE", "OWNS", "INTERESTED_IN", "RELATED_TO"',
        "special_rule": "Capture personal details and investment preferences."
    }
}

# Template Master
MASTER_PROMPT_TEMPLATE = """
ROLE: {role}
TASK: Analyze the provided text chunk and extract a Knowledge Graph in JSON format.

--- CONTEXT INFORMATION ---
SOURCE FILENAME: "{filename}"
PROCESSING DATE: {today}
TEMPORAL HINT: {temporal_hint}

--- RIGID ONTOLOGY (Use ONLY these types in English) ---
NODES: {nodes}
CATCH-ALL NODE: "KeyTerm" (Use for critical concepts that don't fit above).
RELATIONS: {edges}
CATCH-ALL RELATION: "RELATED_TO".

--- RULES ---
1. {special_rule}
2. Output must be strictly VALID JSON.
3. Normalize all entity types and relation names to ENGLISH.
4. Keep specific proper names (e.g., "Mario Rossi", "BTP Italia") in original language.

--- EXAMPLE OUTPUT ---
{{"nodes": [{{"id": "EntityName", "type": "NodeType"}}], "edges": [{{"source": "EntityName", "target": "OtherEntity", "relation": "RELATION_TYPE"}}]}}
"""

# --- INIZIALIZZAZIONE ---
print("🔌 Connessione servizi...")
embedder = SentenceTransformer(MODEL_EMBEDDING_ID, device='cuda')
qdrant_client = QdrantClient(host=QDRANT_HOST, port=6333)
neo4j_driver = GraphDatabase.driver(NEO4J_URI, auth=NEO4J_AUTH)
llm_client = OpenAI(base_url=LM_STUDIO_URL, api_key=LM_STUDIO_API_KEY)

# --- UTILS ---
def calculate_hash(data_bytes):
    return hashlib.md5(data_bytes).hexdigest()

def build_dynamic_prompt(doc_type: str, filename: str) -> str:
    """Costruisce il prompt 'su misura' per il file specifico."""
    
    # 1. Recupera l'ontologia base
    ontology = ONTOLOGIES.get(doc_type, ONTOLOGIES["financial"])
    
    # 2. Cerca indizi temporali nel nome file (es. "Report_2023.pdf")
    # Cerca un anno tra il 1990 e il 2030
    match = re.search(r"(19|20)\d{2}", filename)
    if match:
        year = match.group(0)
        temporal_hint = f"The filename contains '{year}'. Assume unstated dates refer to {year}."
    else:
        temporal_hint = "No specific year in filename. Infer dates from text if possible."

    # 3. Formatta il template
    return MASTER_PROMPT_TEMPLATE.format(
        role=ontology["role"],
        filename=filename,
        today=datetime.datetime.now().strftime("%Y-%m-%d"),
        temporal_hint=temporal_hint,
        nodes=ontology["nodes"],
        edges=ontology["edges"],
        special_rule=ontology["special_rule"]
    )

# --- POSTGRESQL HANDLERS ---
def pg_start_log(filename: str, source_type: str) -> int:
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    try:
        cur.execute("""
            INSERT INTO ingestion_logs (source_name, source_type, status, ingestion_ts)
            VALUES (%s, %s, 'PROCESSING', NOW())
            RETURNING log_id
        """, (filename, source_type))
        log_id = cur.fetchone()[0]
        conn.commit()
        return log_id
    except Exception as e:
        conn.rollback()
        raise e
    finally:
        cur.close()
        conn.close()

def pg_insert_chunk(log_id: int, chunk_index: int, toon_type: str, content: str, metadata: dict):
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    try:
        cur.execute("""
            INSERT INTO document_chunks (log_id, chunk_index, toon_type, content_raw, content_semantic, metadata_json)
            VALUES (%s, %s, %s, %s, %s, %s)
        """, (log_id, chunk_index, toon_type, content, content, Json(metadata)))
        conn.commit()
    except Exception as e:
        conn.rollback()
        print(f"⚠️ Errore insert chunk PG: {e}")
    finally:
        cur.close()
        conn.close()

def pg_save_image(log_id: int, image_bytes: bytes, mime_type: str, description: str) -> int:
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    img_hash = calculate_hash(image_bytes)
    try:
        cur.execute("SELECT image_id FROM ingestion_images WHERE image_hash = %s", (img_hash,))
        res = cur.fetchone()
        if res: return res[0]
        
        cur.execute("""
            INSERT INTO ingestion_images (log_id, image_data, image_hash, mime_type, description_ai, ingestion_ts)
            VALUES (%s, %s, %s, %s, %s, NOW())
            RETURNING image_id
        """, (log_id, psycopg2.Binary(image_bytes), img_hash, mime_type, description))
        image_id = cur.fetchone()[0]
        conn.commit()
        return image_id
    except Exception as e:
        conn.rollback()
        print(f"⚠️ Errore insert image PG: {e}")
        return -1
    finally:
        cur.close()
        conn.close()

def pg_close_log(log_id: int, status: str, total_chunks: int, processing_ms: int, error_msg: str = None):
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    try:
        cur.execute("""
            UPDATE ingestion_logs 
            SET status = %s, total_chunks = %s, processing_time_ms = %s, error_message = %s
            WHERE log_id = %s
        """, (status, total_chunks, processing_ms, error_msg, log_id))
        conn.commit()
    except Exception as e:
        print(f"⚠️ Errore chiusura log PG: {e}")
    finally:
        cur.close()
        conn.close()

# --- AI ENGINES ---
def analyze_image_with_vision(image_bytes: bytes) -> str:
    try:
        base64_image = base64.b64encode(image_bytes).decode('utf-8')
        response = llm_client.chat.completions.create(
            model=MODEL_VISION_ID,
            messages=[
                {"role": "system", "content": "You are a financial analyst."},
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe this image in ENGLISH. Extract key data points, axis labels, and trends."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                    ],
                }
            ],
            temperature=0.1,
            max_tokens=300
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"   ⚠️ Vision AI Error: {e}")
        return "Description not available."

def extract_graph_data(text: str, doc_type: str, filename: str) -> Dict:
    # --- CHIAMATA DINAMICA ---
    # Costruiamo il prompt specifico per questo file in questo momento
    system_prompt = build_dynamic_prompt(doc_type, filename)

    try:
        response = llm_client.chat.completions.create(
            model=MODEL_LLM_ID,
            messages=[
                {"role": "system", "content": system_prompt}, 
                {"role": "user", "content": f"TEXT TO ANALYZE:\n{text}"}
            ],
            temperature=0.0
        )
        content = response.choices[0].message.content.strip()
        if "```json" in content: content = content.split("```json")[1].split("```")[0].strip()
        elif "```" in content: content = content.split("```")[1].split("```")[0].strip()
        return json.loads(content)
    except Exception as e:
        # print(f"Errore JSON LLM: {e}") # Decommentare per debug
        return None

def save_to_neo4j(chunk_uuid: str, graph_data: Dict, doc_type: str):
    if not graph_data: return
    
    with neo4j_driver.session() as session:
        # 1. Crea il Chunk Node (questo è sicuro)
        session.run("MERGE (c:Chunk {id: $cid, type: $dtype})", cid=chunk_uuid, dtype=doc_type)
        
        # 2. Processa Nodi (con protezione)
        for node in graph_data.get("nodes", []):
            # Controllo esistenza chiavi obbligatorie
            if "id" not in node or "type" not in node:
                continue # Salta nodo malformato
                
            try:
                session.run(
                    "MERGE (n:Entity {name: $name}) ON CREATE SET n.type = $type "
                    "WITH n MATCH (c:Chunk {id: $cid}) MERGE (n)-[:MENTIONED_IN]->(c)", 
                    name=node['id'], type=node['type'], cid=chunk_uuid
                )
            except Exception as e:
                print(f"   ⚠️ Neo4j Node Error: {e}")

        # 3. Processa Archi (con protezione)
        for edge in graph_data.get("edges", []):
            # Controllo esistenza chiavi obbligatorie
            if "source" not in edge or "target" not in edge or "relation" not in edge:
                continue # Salta arco malformato
            
            try:
                rel_type = edge['relation'].upper().replace(" ", "_")
                # Sanitize relation name (solo lettere, numeri e underscore)
                rel_type = "".join(c for c in rel_type if c.isalnum() or c == '_')
                
                if not rel_type: rel_type = "RELATED_TO" # Fallback

                session.run(
                    f"MATCH (a:Entity {{name: $s}}), (b:Entity {{name: $t}}) MERGE (a)-[:{rel_type}]->(b)",
                    s=edge['source'], t=edge['target']
                )
            except Exception as e:
                print(f"   ⚠️ Neo4j Edge Error: {e}")
# --- FILE PARSERS ---
def extract_pdf_content(file_path: str, log_id: int) -> str:
    full_text = ""
    doc = fitz.open(file_path)
    print(f"   🔍 Analisi PDF ({len(doc)} pagine)...")
    for page in doc:
        full_text += page.get_text() + "\n"
        image_list = page.get_images(full=True)
        for img_info in image_list:
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            ext = base_image["ext"]
            if len(image_bytes) < 3000: continue
            
            print("   👁️ Vision AI processing...")
            description = analyze_image_with_vision(image_bytes)
            img_id = pg_save_image(log_id, image_bytes, f"image/{ext}", description)
            full_text += f"\n\n[[IMAGE_DESC_AI (ID: {img_id}): {description}]]\n\n"
    doc.close()
    return full_text

def get_file_content(file_path: str, log_id: int) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf": return extract_pdf_content(file_path, log_id)
    elif ext in [".docx", ".doc"]: return "\n".join([p.text for p in Document(file_path).paragraphs])
    elif ext in [".pptx", ".ppt"]:
        text = []
        for slide in Presentation(file_path).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text: text.append(shape.text)
        return "\n".join(text)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f: return f.read()
    return ""

# --- MAIN LOOP ---
def process_single_file(file_path: str, doc_type: str):
    filename = os.path.basename(file_path)
    start_time = time.time()
    print(f"\n📄 Start: {filename} ({doc_type})")
    
    try:
        log_id = pg_start_log(filename, doc_type)
    except Exception as e:
        print(f"❌ DB Init Error: {e}")
        return

    try:
        content = get_file_content(file_path, log_id)
        if not content.strip():
            pg_close_log(log_id, "SKIPPED_EMPTY", 0, int((time.time()-start_time)*1000))
            return

        chunk_size = 1000
        chunks_text = [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]
        
        count = 0
        for i, text_chunk in enumerate(chunks_text):
            if len(text_chunk) < 50: continue
            
            chunk_uuid = str(uuid.uuid4())
            toon_type = "image_description" if "[[IMAGE_DESC_AI" in text_chunk else "text"
            
            vec = embedder.encode(text_chunk, normalize_embeddings=True).tolist()
            qdrant_client.upsert(
                collection_name=QDRANT_COLLECTION,
                points=[models.PointStruct(
                    id=chunk_uuid, 
                    vector=vec, 
                    payload={"filename": filename, "type": doc_type, "chunk_index": i, "toon_type": toon_type}
                )]
            )
            
            # CHIAMATA AGGIORNATA: Passiamo anche il filename per il prompt dinamico
            g_data = extract_graph_data(text_chunk, doc_type, filename)
            
            if g_data: save_to_neo4j(chunk_uuid, g_data, doc_type)
            
            meta = {"qdrant_uuid": chunk_uuid, "original_path": file_path}
            pg_insert_chunk(log_id, i, toon_type, text_chunk, meta)
            
            print(f"   chunk {i} ({toon_type}): Vector ✅ | Graph {'✅' if g_data else '❌'}")
            count += 1

        elapsed = int((time.time() - start_time) * 1000)
        pg_close_log(log_id, "COMPLETED", count, elapsed)
        archive_file(file_path, doc_type)

    except Exception as e:
        elapsed = int((time.time() - start_time) * 1000)
        print(f"❌ Critical Error: {e}")
        pg_close_log(log_id, "FAILED", 0, elapsed, str(e))

def archive_file(file_path: str, category: str):
    try:
        dest = os.path.join(PROCESSED_DIR, category)
        os.makedirs(dest, exist_ok=True)
        fname = os.path.basename(file_path)
        name, ext = os.path.splitext(fname)
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        shutil.move(file_path, os.path.join(dest, f"{name}_{ts}{ext}"))
        print(f"   📦 Archived.")
    except Exception as e:
        print(f"   ❌ Archive Error: {e}")

def main():
    if not qdrant_client.collection_exists(QDRANT_COLLECTION):
        qdrant_client.create_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=models.VectorParams(size=1024, distance=models.Distance.COSINE)
        )

    folder_map = {
        "financial_reports": "financial",
        "educational": "educational",
        "client_profiles": "client"
    }
    exts = ["*.txt", "*.pdf", "*.docx", "*.doc", "*.pptx"]

    for folder, doc_type in folder_map.items():
        base = os.path.join(BASE_DATA_DIR, folder)
        files = []
        for ext in exts: files.extend(glob.glob(os.path.join(base, ext)))
        
        if files:
            print(f"📂 Folder '{folder}': {len(files)} files.")
            for f in files: process_single_file(f, doc_type)

    print("\n✅ Ingestion Finished.")
    neo4j_driver.close()

if __name__ == "__main__":
    for d in ["financial_reports", "educational", "client_profiles"]:
        os.makedirs(os.path.join(BASE_DATA_DIR, d), exist_ok=True)
    os.makedirs(PROCESSED_DIR, exist_ok=True)
    main()