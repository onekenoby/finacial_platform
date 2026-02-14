"""
Ingestion Engine (Production-oriented)
- Supporto: PDF, DOCX, PPTX, TXT
- DB: PostgreSQL (TimescaleDB), Qdrant, Neo4j
- AI locale (LM Studio): LLM (KG + classifier) + Vision (immagini PDF)
- Embeddings locali: SentenceTransformers (batch su GPU)
- TOON: gestione chunk + metadata + (opzionale) semantic normalization

Migliorie applicate (rispetto al tuo codice originale):
1) Prompt più robusti (classifier + KG + vision) con vincoli anti-hallucination
2) Embeddings calcolati in batch (GPU sfruttata meglio)
3) Riduzione drastica chiamate LLM (KG solo su chunk candidati + limite max per documento)
4) Chunking migliore (PDF page-aware + split semantico con overlap; DOCX per paragrafi; PPTX per shape)
5) Neo4j write path ottimizzato (UNWIND batch; no query per nodo/arco)
   - Relationship uniforme :REL con property "type" per evitare rel-type dinamici (più sicuro e più batch-friendly)
6) Postgres connection pooling (meno overhead, più velocità e stabilità)

ULTERIORI OTTIMIZZAZIONI AGGIUNTE (questa versione):
A) Postgres immagini: riuso conn/cursor per TUTTO il PDF (evita migliaia di acquire/release dal pool)
B) Neo4j: flush in UNICA chiamata (UNWIND su lista jobs), non più 1 call per chunk
C) Facts extraction: segnali deterministici (%, bps, valute, date) salvati in metadata_json per chunk
"""

import os
import re
import json
import time
import uuid
import shutil
import hashlib
import datetime
import base64 
from typing import List, Dict, Tuple, Optional

import fitz  # PyMuPDF
import psycopg2
from psycopg2.extras import Json, execute_values
from psycopg2.pool import SimpleConnectionPool

from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient, models
from neo4j import GraphDatabase
from openai import OpenAI

# =========================
# CONFIGURAZIONE
# =========================
BASE_DATA_DIR = "./data_ingestion"
INBOX_DIR = os.path.join(BASE_DATA_DIR, "INBOX")
PROCESSED_DIR = os.path.join(BASE_DATA_DIR, "PROCESSED")

CATEGORIES = {
    "pdf": ["pdf"],
    "docx": ["docx", "doc"],
    "pptx": ["pptx", "ppt"],
    "txt": ["txt"]
}

# Chunking
CHUNK_MAX_CHARS = 1400
CHUNK_OVERLAP_CHARS = 200
MIN_CHUNK_LEN = 40

# Embedding batching
DB_FLUSH_SIZE = 256
EMBED_BATCH_SIZE = 64

# KG extraction gating
MAX_KG_CHUNKS_PER_DOC = 10
KG_MIN_LEN = 400
KG_KEYWORDS = ["risk", "guidance", "forecast", "revenue", "earnings", "inflation", "rate", "spread", "debt", "cash"]

# PDF Vision extraction
PDF_VISION_ENABLED = True
PDF_VISION_ONLY_IF_TEXT_SCARSO = True
PDF_MIN_TEXT_LEN_FOR_NO_VISION = 900
PDF_MAX_IMAGES_PER_PAGE = 2
PDF_MIN_IMAGE_BYTES = 35_000
PDF_MIN_IMAGE_DIM = 240

# Qdrant
QDRANT_HOST = os.getenv("QDRANT_HOST", "127.0.0.1")
QDRANT_PORT = int(os.getenv("QDRANT_PORT", "6333"))
QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", "financial_docs")

# Postgres
PG_HOST = os.getenv("PG_HOST", "127.0.0.1")
PG_PORT = int(os.getenv("PG_PORT", "5432"))
PG_DB = os.getenv("PG_DB", "ai_ingestion")
PG_USER = os.getenv("PG_USER", "admin")
PG_PASS = os.getenv("PG_PASS", "admin_password")
PG_MIN_CONN = int(os.getenv("PG_MIN_CONN", "1"))
PG_MAX_CONN = int(os.getenv("PG_MAX_CONN", "5"))

# Neo4j
NEO4J_URI = os.getenv("NEO4J_URI", "bolt://127.0.0.1:7687")
NEO4J_USER = os.getenv("NEO4J_USER", "neo4j")
NEO4J_PASS = os.getenv("NEO4J_PASS", "password_sicura")

# LM Studio / OpenAI-compatible endpoint
LM_BASE_URL = os.getenv("LM_BASE_URL", "http://127.0.0.1:1234/v1")
LM_API_KEY = os.getenv("LM_API_KEY", "lm-studio")

LLM_MODEL_NAME = os.getenv("LLM_MODEL_NAME", "local-model")
VISION_MODEL_NAME = os.getenv("VISION_MODEL_NAME", "local-vision-model")
EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "sentence-transformers/all-MiniLM-L6-v2")

# =========================
# CLIENTS INIT
# =========================
openai_client = OpenAI(base_url=LM_BASE_URL, api_key=LM_API_KEY)

embedder = SentenceTransformer(EMBEDDING_MODEL_NAME)

qdrant_client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)

neo4j_driver = GraphDatabase.driver(NEO4J_URI, auth=(NEO4J_USER, NEO4J_PASS))

pg_pool = SimpleConnectionPool(
    PG_MIN_CONN,
    PG_MAX_CONN,
    host=PG_HOST,
    port=PG_PORT,
    dbname=PG_DB,
    user=PG_USER,
    password=PG_PASS
)

# =========================
# PROMPTS
# =========================
CLASSIFIER_PROMPT = """
You are a strict classifier.
Given a chunk of text, decide if it is likely useful to extract entities/relations for a knowledge graph.
Return ONLY valid JSON:
{"useful": true/false, "reason": "...", "signals": ["..."]}
Be conservative: if uncertain, return useful=false.
"""

KG_PROMPT = """
Extract entities and relations from the text for a knowledge graph.

Return ONLY valid JSON with this schema:
{
  "nodes": [{"id": "...", "label": "...", "type": "..."}],
  "edges": [{"source": "...", "target": "...", "relation": "...", "label": "..."}]
}

Rules:
- Do not invent facts. If unsure, omit.
- Keep nodes and edges concise.
- Use stable ids (normalized strings), avoid long ids.
"""

VISION_PROMPT = """
You analyze an image extracted from a PDF page and provide useful structured information.

Return ONLY valid JSON:
{
  "kind": "table|chart|diagram|photo|other",
  "key_points": ["..."],   // max 6
  "numbers": [{"label":"", "value":"", "unit":"", "currency":"", "period":""}],  // only if clearly readable
  "unreadable_parts": ["..."]  // optional
}

If the image is decorative or unreadable, use kind="other" and empty arrays.
"""

# =========================
# UTILS
# =========================
def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()

def normalize_ws(text: str) -> str:
    # Normalizzazione semplice ma utile: whitespace + line breaks
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def extract_facts(text: str) -> Dict:
    """
    Estrae segnali 'hard' (veloci e deterministici) dal testo: numeri, valute, percentuali, date.
    Utile per KB finanziaria: riduce rumore e aiuta re-ranking / answer grounding.
    """
    if not text:
        return {}

    # Limita testo per performance (evita regex su chunk enormi)
    t = text[:20000]

    # Percentuali / bps
    perc = re.findall(r"\b\d+(?:[\.,]\d+)?\s?%\b", t)
    bps = re.findall(r"\b\d+(?:[\.,]\d+)?\s?bps\b", t, flags=re.IGNORECASE)

    # Valute e importi (pattern conservativi)
    currency = re.findall(r"(?:€\s?\d[\d\.,]*|\$\s?\d[\d\.,]*|£\s?\d[\d\.,]*|\b\d[\d\.,]*\s?(?:EUR|USD|GBP)\b)", t)

    # Date comuni
    dates = re.findall(r"\b\d{4}-\d{2}-\d{2}\b|\b\d{1,2}/\d{1,2}/\d{2,4}\b", t)

    # Anni (solo se appaiono vicino a parole chiave tipiche, per evitare rumore)
    years = []
    for m in re.finditer(r"\b(19|20)\d{2}\b", t):
        s = t[max(0, m.start()-24): m.end()+24].lower()
        if any(k in s for k in ["fy", "q", "quarter", "anno", "year", "202", "guidance", "forecast", "outlook", "cagr"]):
            years.append(m.group(0))

    def _uniq_limit(xs, limit=30):
        out=[]
        seen=set()
        for x in xs:
            x=x.strip()
            if not x or x in seen:
                continue
            seen.add(x)
            out.append(x)
            if len(out)>=limit:
                break
        return out

    facts = {}
    if perc:
        facts["percentages"] = _uniq_limit(perc)
    if bps:
        facts["bps"] = _uniq_limit(bps)
    if currency:
        facts["amounts"] = _uniq_limit(currency)
    if dates:
        facts["dates"] = _uniq_limit(dates)
    if years:
        facts["years_contextual"] = _uniq_limit(years, limit=20)

    return facts

def normalize_entity_id(s: str) -> str:
    s = (s or "").strip().lower()
    s = s.replace('"', "").replace("'", "")
    s = re.sub(r"\s+", " ", s)
    return s[:180]  # evita id troppo lunghi

def safe_json_extract(raw: str) -> Optional[Dict]:
    try:
        raw = raw.strip()
        # Rimuovi eventuali code fence
        raw = re.sub(r"^```json", "", raw, flags=re.IGNORECASE).strip()
        raw = re.sub(r"```$", "", raw).strip()
        return json.loads(raw)
    except Exception:
        return None

def split_text_with_overlap(text: str, max_chars: int, overlap: int) -> List[str]:
    text = normalize_ws(text)
    if len(text) <= max_chars:
        return [text]
    out = []
    i = 0
    while i < len(text):
        chunk = text[i:i+max_chars]
        out.append(chunk)
        i += (max_chars - overlap)
        if i < 0:
            break
    return out

def split_paragraphs(text: str, max_chars: int, overlap: int) -> List[str]:
    # Split semplice per paragrafi, poi fallback su overlap
    text = normalize_ws(text)
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    merged = []
    buf = ""
    for p in paras:
        if len(buf) + len(p) + 2 <= max_chars:
            buf = (buf + "\n\n" + p).strip()
        else:
            if buf:
                merged.append(buf)
            buf = p
    if buf:
        merged.append(buf)

    # Ora applica overlap sul testo unito
    final_chunks = []
    for m in merged:
        final_chunks.extend(split_text_with_overlap(m, max_chars, overlap))
    return final_chunks

def is_candidate_for_kg(text: str, doc_type: str) -> bool:
    t = text.lower()
    if len(t) < KG_MIN_LEN:
        return False
    if any(k in t for k in KG_KEYWORDS):
        return True
    return False

# =========================
# POSTGRES (POOL)
# =========================
def pg_get_conn():
    return pg_pool.getconn()

def pg_put_conn(conn):
    pg_pool.putconn(conn)

def pg_start_log(source_name: str, source_type: str) -> int:
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO ingestion_logs (source_name, source_type, ingestion_ts, status) VALUES (%s, %s, NOW(), %s) RETURNING log_id",
                (source_name, source_type, "RUNNING")
            )
            log_id = cur.fetchone()[0]
        conn.commit()
        return log_id
    finally:
        pg_put_conn(conn)

def pg_close_log(log_id: int, status: str, total_chunks: int, processing_ms: int, error_msg: str = None):
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE ingestion_logs SET status = %s, total_chunks = %s, processing_time_ms = %s, error_message = %s "
                "WHERE log_id = %s",
                (status, total_chunks, processing_ms, error_msg, log_id)
            )
        conn.commit()
    except Exception:
        conn.rollback()
    finally:
        pg_put_conn(conn)

def pg_get_image_by_hash(image_hash: str, cur=None) -> Optional[Tuple[int, str]]:
    """Ritorna (image_id, description_ai) se esiste. Se cur è fornito, riusa la stessa transazione."""
    if cur is not None:
        cur.execute(
            "SELECT image_id, description_ai FROM ingestion_images WHERE image_hash = %s LIMIT 1",
            (image_hash,)
        )
        return cur.fetchone()

    conn = pg_get_conn()
    try:
        with conn.cursor() as _cur:
            _cur.execute(
                "SELECT image_id, description_ai FROM ingestion_images WHERE image_hash = %s LIMIT 1",
                (image_hash,)
            )
            return _cur.fetchone()
    except Exception:
        return None
    finally:
        pg_put_conn(conn)

def pg_save_image(log_id: int, image_bytes: bytes, mime_type: str, description: str, cur=None) -> int:
    """
    Inserisce solo se non esiste già per hash.
    - Se cur è fornito, NON fa commit/rollback (li gestisce il chiamante).
    - Se cur è None, apre una connessione dal pool e committa internamente.
    """
    img_hash = sha256_hex(image_bytes)

    # cache check (stessa transazione se cur è fornito)
    cached = pg_get_image_by_hash(img_hash, cur=cur)
    if cached:
        return cached[0]

    if cur is not None:
        cur.execute(
            "INSERT INTO ingestion_images (log_id, image_data, image_hash, mime_type, description_ai, ingestion_ts) "
            "VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING image_id",
            (log_id, psycopg2.Binary(image_bytes), img_hash, mime_type, description)
        )
        return cur.fetchone()[0]

    conn = pg_get_conn()
    try:
        with conn.cursor() as _cur:
            _cur.execute(
                "INSERT INTO ingestion_images (log_id, image_data, image_hash, mime_type, description_ai, ingestion_ts) "
                "VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING image_id",
                (log_id, psycopg2.Binary(image_bytes), img_hash, mime_type, description)
            )
            image_id = _cur.fetchone()[0]
        conn.commit()
        return image_id
    except Exception:
        conn.rollback()
        return -1
    finally:
        pg_put_conn(conn)

def flush_postgres_chunks_batch(batch_data: List[Tuple]):
    if not batch_data:
        return
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            execute_values(
                cur,
                """
                INSERT INTO document_chunks
                    (log_id, chunk_index, toon_type, content_raw, content_semantic, metadata_json)
                VALUES %s
                """,
                batch_data
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        print(f"   ⚠️ Postgres Batch Error: {e}")
    finally:
        pg_put_conn(conn)

# =========================
# NEO4J BATCH (UNWIND)
# =========================
NEO4J_QUERY = """
// 1) Chunk node
MERGE (c:Chunk {id: $chunk_id})
SET c.doc_type = $doc_type,
    c.filename = $filename,
    c.log_id = $log_id,
    c.chunk_index = $chunk_index,
    c.toon_type = $toon_type,
    c.page_no = $page_no,
    c.ingested_at = datetime()

// 2) Entities + mention relationship
WITH c
UNWIND $nodes AS n
MERGE (e:Entity {id: n.id})
ON CREATE SET e.label = n.label, e.type = n.type
ON MATCH  SET e.label = coalesce(e.label, n.label),
             e.type  = coalesce(e.type, n.type)
MERGE (c)-[:MENTIONS]->(e)

// 3) Relations
WITH $edges AS edges
UNWIND edges AS ed
MERGE (s:Entity {id: ed.source})
MERGE (t:Entity {id: ed.target})
MERGE (s)-[rel:REL {type: ed.relation}]->(t)
SET rel.label = coalesce(ed.label, ed.relation),
    rel.log_id = $log_id,
    rel.filename = $filename
"""

NEO4J_BATCH_QUERY = """
UNWIND $rows AS r

// 1. Crea/Trova il Nodo Chunk (Ancoraggio)
MERGE (c:Chunk {id: r.chunk_id})
SET c.filename = r.filename, 
    c.page = r.page_no, 
    c.type = r.doc_type,
    c.ingested_at = datetime()

// 2. Crea i Nodi Entità con Label Dinamica
WITH c, r
UNWIND r.nodes AS n
// APOC permette di passare la label come variabile (n.type)
// Esempio: Se n.type è "Company", crea (e:Company:Entity {id: "..."})
CALL apoc.merge.node([n.type, "Entity"], {id: n.id}, {label: n.label}, {}) YIELD node as e
// Collega l'entità al Chunk (Provenance)
MERGE (e)-[:MENTIONED_IN]->(c)

// 3. Crea le Relazioni con Tipo Dinamico
WITH r
UNWIND r.edges AS rel
MATCH (s {id: rel.source}) 
MATCH (t {id: rel.target})
// APOC per creare la relazione dinamica (es. :OWNS invece di :REL)
CALL apoc.merge.relationship(s, rel.relation, {}, {}, t, {}) YIELD rel as r_out
RETURN count(*)
"""

def _sanitize_graph(graph: Dict) -> Tuple[List[Dict], List[Dict]]:
    """Dedup, normalizza e prepara per APOC (tipi dinamici)."""
    if not graph: 
        return [], []
    
    raw_nodes = graph.get("nodes", []) or []
    raw_edges = graph.get("edges", []) or []
    
    # --- NODI ---
    dedup_nodes = []
    seen_ids = set()

    for n in raw_nodes:
        # 1. Normalizza ID
        nid = normalize_entity_id(n.get("id"))
        if not nid or nid in seen_ids: 
            continue
        
        # 2. Pulisci il Tipo (PascalCase per Neo4j Labels)
        # Es: "financial institution" -> "FinancialInstitution"
        raw_type = n.get("type", "Entity").strip()
        # Rimuove caratteri speciali e mette in maiuscolo le prime lettere
        clean_type = "".join(x.capitalize() for x in re.split(r"[_\s\-\.]+", raw_type) if x)
        # Fallback se esce stringa vuota
        ntype = re.sub(r"[^a-zA-Z0-9]", "", clean_type) or "Entity"

        dedup_nodes.append({
            "id": nid,
            "label": (n.get("label") or nid)[:200], # Label umana
            "type": ntype  # Label tecnica per Neo4j (es. :Company)
        })
        seen_ids.add(nid)

    # --- ARCHI ---
    dedup_edges = []
    seen_edges = set()

    for e in raw_edges:
        src = normalize_entity_id(e.get("source"))
        tgt = normalize_entity_id(e.get("target"))
        
        if not src or not tgt:
            continue

        # 3. Pulisci la Relazione (UPPER_SNAKE_CASE per Neo4j Relationships)
        # Es: "owns shares in" -> "OWNS_SHARES_IN"
        raw_rel = e.get("relation", "RELATED_TO").upper()
        # Sostituisce spazi e trattini con underscore, rimuove altro
        rel = re.sub(r"[^A-Z0-9_]+", "_", raw_rel).strip("_")
        if not rel: 
            rel = "RELATED_TO"

        # Chiave univoca per evitare duplicati nello stesso batch
        key = (src, tgt, rel)
        if key in seen_edges: 
            continue
        
        dedup_edges.append({
            "source": src,
            "target": tgt,
            "relation": rel
        })
        seen_edges.add(key)

    return dedup_nodes, dedup_edges
    """Dedup e normalizza nodes/edges per Neo4j."""
    if not graph:
        return [], []

    raw_nodes = graph.get("nodes", []) or []
    raw_edges = graph.get("edges", []) or []

    clean_nodes = []
    for n in raw_nodes:
        if not isinstance(n, dict):
            continue
        nid = normalize_entity_id(n.get("id") or n.get("label") or "")
        if not nid:
            continue
        clean_nodes.append({
            "id": nid,
            "label": (n.get("label") or nid)[:300],
            "type": (n.get("type") or "")[:80]
        })

    clean_edges = []
    for e in raw_edges:
        if not isinstance(e, dict):
            continue
        src = normalize_entity_id(e.get("source") or "")
        tgt = normalize_entity_id(e.get("target") or "")
        rel = (e.get("relation") or e.get("type") or "").strip()
        rel = rel[:80] if rel else "RELATED_TO"
        if not src or not tgt:
            continue
        clean_edges.append({
            "source": src,
            "target": tgt,
            "relation": rel,
            "label": (e.get("label") or rel)[:200]
        })

    # Dedup
    seen_n = set()
    dedup_nodes = []
    for n in clean_nodes:
        if n["id"] in seen_n:
            continue
        seen_n.add(n["id"])
        dedup_nodes.append(n)

    seen_e = set()
    dedup_edges = []
    for e in clean_edges:
        key = (e["source"], e["target"], e["relation"])
        if key in seen_e:
            continue
        seen_e.add(key)
        dedup_edges.append(e)

    return dedup_nodes, dedup_edges

def flush_neo4j_jobs_batch(jobs: List[Tuple], log_id: int, doc_type: str, filename: str):
    """
    Flush Neo4j in un'unica chiamata UNWIND (molto più veloce di 1 chiamata per chunk).
    jobs: (chunk_uuid, chunk_index, toon_type, page_no, graph)
    """
    if not jobs:
        return
    rows = []
    for job in jobs:
        try:
            cid, cidx, ttype, pno, g = job
            nodes, edges = _sanitize_graph(g or {})
            if not nodes and not edges:
                continue
            rows.append({
                "chunk_id": cid,
                "log_id": log_id,
                "doc_type": doc_type,
                "filename": filename,
                "chunk_index": int(cidx),
                "toon_type": ttype,
                "page_no": int(pno) if pno else None,
                "nodes": nodes,
                "edges": edges
            })
        except Exception:
            continue

    if not rows:
        return

    with neo4j_driver.session() as session:
        session.run(NEO4J_BATCH_QUERY, rows=rows)

def flush_neo4j_chunk_graph(chunk_id: str, log_id: int, doc_type: str, filename: str,
                           chunk_index: int, toon_type: str, page_no: int, graph: Dict):
    if not graph:
        return
    nodes = graph.get("nodes", []) or []
    edges = graph.get("edges", []) or []

    # Normalizza + dedup
    clean_nodes = []
    for n in nodes:
        nid = normalize_entity_id(n.get("id") or n.get("label") or "")
        if not nid:
            continue
        clean_nodes.append({
            "id": nid,
            "label": (n.get("label") or nid)[:300],
            "type": (n.get("type") or "")[:80]
        })

    clean_edges = []
    for e in edges:
        src = normalize_entity_id(e.get("source") or "")
        tgt = normalize_entity_id(e.get("target") or "")
        rel = (e.get("relation") or e.get("type") or "").strip()
        rel = rel[:80] if rel else "RELATED_TO"
        if not src or not tgt:
            continue
        clean_edges.append({
            "source": src,
            "target": tgt,
            "relation": rel,
            "label": (e.get("label") or rel)[:200]
        })

    seen_n = set()
    dedup_nodes = []
    for n in clean_nodes:
        if n["id"] in seen_n:
            continue
        seen_n.add(n["id"])
        dedup_nodes.append(n)

    seen_e = set()
    dedup_edges = []
    for e in clean_edges:
        key = (e["source"], e["target"], e["relation"])
        if key in seen_e:
            continue
        seen_e.add(key)
        dedup_edges.append(e)

    if not dedup_nodes and not dedup_edges:
        return

    with neo4j_driver.session() as session:
        try:
            session.run(
                NEO4J_QUERY,
                chunk_id=chunk_id,
                log_id=log_id,
                doc_type=doc_type,
                filename=filename,
                chunk_index=chunk_index,
                toon_type=toon_type,
                page_no=page_no,
                nodes=dedup_nodes,
                edges=dedup_edges
            )
        except Exception as e:
            print(f"   ⚠️ Neo4j Error: {e}")

# =========================
# LLM CALLS (LM Studio)
# =========================
def llm_chat(prompt: str, user_text: str, model: str) -> str:
    resp = openai_client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": prompt.strip()},
            {"role": "user", "content": user_text.strip()}
        ],
        temperature=0.1,
        max_tokens=900
    )
    return resp.choices[0].message.content or ""

def classify_chunk_useful(text: str) -> bool:
    out = llm_chat(CLASSIFIER_PROMPT, text, LLM_MODEL_NAME)
    js = safe_json_extract(out) or {}
    return bool(js.get("useful"))

def extract_graph_data(text: str, doc_type: str, filename: str) -> Optional[Dict]:
    # Extra safety: classification step can be enabled; currently we use heuristic gating + budget
    # If you want even less KG calls, uncomment:
    # if not classify_chunk_useful(text):
    #     return None
    out = llm_chat(KG_PROMPT, text, LLM_MODEL_NAME)
    js = safe_json_extract(out)
    return js



def analyze_image_with_vision(image_bytes: bytes) -> Optional[Dict]:
    if not image_bytes:
        return None
        
    try:
        # 1. Encode base64
        b64_image = base64.b64encode(image_bytes).decode('utf-8')
        
        # 2. Call Vision Model
        response = openai_client.chat.completions.create(
            model=VISION_MODEL_NAME, # Assicurati che punti a Qwen/Llava in LM Studio
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": VISION_PROMPT},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{b64_image}"
                            },
                        },
                    ],
                }
            ],
            max_tokens=500,
            temperature=0.1
        )
        
        # 3. Extract JSON
        content = response.choices[0].message.content
        return safe_json_extract(content)

    except Exception as e:
        print(f" ⚠️ Vision Error: {e}")
        return None

# =========================
# EXTRACTION: DOCX/PPTX/TXT
# =========================
def extract_docx_chunks(file_path: str) -> List[Dict]:
    try:
        import docx
    except ImportError:
        return []
    doc = docx.Document(file_path)
    paras = [normalize_ws(p.text) for p in doc.paragraphs if p.text and p.text.strip()]
    out = []
    for p in paras:
        for c in split_text_with_overlap(p, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
            out.append({"text_raw": c, "text_sem": c, "page_no": None, "toon_type": "text"})
    return out

def extract_pptx_chunks(file_path: str) -> List[Dict]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                buf.append(shape.text.strip())
        s = normalize_ws("\n".join(buf))
        if s:
            slides.append((i + 1, s))

    out = []
    for slide_no, s in slides:
        chunks = split_text_with_overlap(s, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)
        for c in chunks:
            out.append({"text_raw": c, "text_sem": c, "page_no": slide_no, "toon_type": "text"})
    return out

def extract_txt_chunks(file_path: str) -> List[Dict]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = normalize_ws(f.read())
    out = []
    for c in split_text_with_overlap(text, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
        out.append({"text_raw": c, "text_sem": c, "page_no": None, "toon_type": "text"})
    return out

# =========================
# PDF UTIL: rimozione righe ripetute (header/footer)
# =========================
def remove_repeated_lines(pages_text: List[str], min_occurrences: int = 3) -> List[str]:
    # conta righe
    freq = {}
    pages_lines = []
    for t in pages_text:
        lines = [ln.strip() for ln in t.splitlines() if ln.strip()]
        pages_lines.append(lines)
        for ln in set(lines):
            freq[ln] = freq.get(ln, 0) + 1

    repeated = {ln for ln, c in freq.items() if c >= min_occurrences and len(ln) <= 120}

    cleaned = []
    for t in pages_text:
        lines = [ln for ln in t.splitlines() if ln.strip() not in repeated]
        cleaned.append("\n".join(lines))
    return cleaned

def extract_pdf_chunks(file_path: str, log_id: int) -> List[Dict]:
    """
    Estrae chunk da PDF:
    - 1 livello: testo per pagina (con pulizia righe ripetute)
    - 2 livello: split semantico per paragrafi + overlap
    - Vision: solo immagini grandi + limitate; caching su hash; opzionale solo se testo pagina scarso
    Restituisce lista di dict: {text_raw, text_sem, page_no, toon_type, images_meta...}
    """
    doc = fitz.open(file_path)
    # Riusa la stessa connessione/cursor Postgres per tutte le immagini del PDF (molto più veloce).
    pg_conn = None
    pg_cur = None

    def _get_pg_cur():
        nonlocal pg_conn, pg_cur
        if pg_conn is None:
            pg_conn = pg_get_conn()
            pg_cur = pg_conn.cursor()
        return pg_cur
    pages_text = [doc[i].get_text("text") for i in range(len(doc))]
    pages_text = [normalize_ws(t) for t in pages_text]
    pages_text = remove_repeated_lines(pages_text)

    out_chunks: List[Dict] = []
    print(f"   🔍 Analisi PDF ({len(doc)} pagine)...")

    for page_no in range(len(doc)):
        page = doc[page_no]
        page_text = pages_text[page_no]
        page_text_len = len(page_text)

        # Vision gating
        allow_vision = True
        if PDF_VISION_ONLY_IF_TEXT_SCARSO:
            allow_vision = page_text_len < PDF_MIN_TEXT_LEN_FOR_NO_VISION

        image_desc_blocks = []
        if PDF_VISION_ENABLED and allow_vision:
            try:
                imgs = page.get_images(full=True) or []
            except Exception:
                imgs = []

            used = 0
            for img in imgs:
                if used >= PDF_MAX_IMAGES_PER_PAGE:
                    break
                try:
                    xref = img[0]
                    base_img = doc.extract_image(xref)
                    img_bytes = base_img.get("image")
                    if not img_bytes or len(img_bytes) < PDF_MIN_IMAGE_BYTES:
                        continue

                    ext = (base_img.get("ext") or "png").lower()

                    # size gating (evita iconcine)
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        if pix.width < PDF_MIN_IMAGE_DIM or pix.height < PDF_MIN_IMAGE_DIM:
                            continue
                    except Exception:
                        pass

                    img_hash = sha256_hex(img_bytes)
                    cached = pg_get_image_by_hash(img_hash, cur=_get_pg_cur())
                    if cached:
                        img_id, desc_ai = cached[0], cached[1] or ""
                    else:
                        # Vision analysis (placeholder, adjust if you have true vision input)
                        vis = analyze_image_with_vision(img_bytes)
                        if vis and isinstance(vis, dict):
                            kind = vis.get("kind", "other")
                            key_points = vis.get("key_points", []) or []
                            numbers = vis.get("numbers", []) or []
                            desc_ai = f"[{kind}] " + "; ".join(key_points[:6])
                            if numbers:
                                desc_ai += " | numbers: " + "; ".join(
                                    [f"{n.get('label','')}: {n.get('value','')}{n.get('unit','')}" for n in numbers[:6]]
                                )
                        else:
                            desc_ai = ""  # se vision non disponibile

                        if not desc_ai:
                            continue
                        img_id = pg_save_image(log_id, img_bytes, f"image/{ext}", desc_ai, cur=_get_pg_cur())

                    if img_id > 0:
                        image_desc_blocks.append(f"[[IMAGE_DESC_AI (ID: {img_id}, PAGE: {page_no+1}): {desc_ai}]]")
                        used += 1
                except Exception:
                    continue

        # Commit immagini della pagina (se abbiamo riusato una connessione)
        if pg_conn is not None:
            try:
                pg_conn.commit()
            except Exception:
                pg_conn.rollback()

        # Costruisci testo per pagina con eventuali blocchi immagini
        combined_raw = page_text
        if image_desc_blocks:
            combined_raw += "\n\n" + "\n\n".join(image_desc_blocks)

        combined_sem = normalize_ws(combined_raw)

        # Split in chunk semantici
        page_chunks = split_paragraphs(combined_sem, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)
        for ch in page_chunks:
            toon_type = "image_description" if "[[IMAGE_DESC_AI" in ch else "text"
            out_chunks.append({
                "text_raw": ch,
                "text_sem": ch,
                "page_no": page_no + 1,
                "toon_type": toon_type
            })

    doc.close()

    # Rilascia risorse PG riusate per immagini
    if pg_cur is not None:
        try:
            pg_cur.close()
        except Exception:
            pass
    if pg_conn is not None:
        pg_put_conn(pg_conn)

    return out_chunks

def extract_file_chunks(file_path: str, log_id: int) -> List[Dict]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_chunks(file_path, log_id)
    elif ext in [".docx", ".doc"]:
        return extract_docx_chunks(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_pptx_chunks(file_path)
    elif ext == ".txt":
        return extract_txt_chunks(file_path)
    return []

# =========================
# QDRANT INIT
# =========================
def ensure_qdrant_collection():
    try:
        qdrant_client.get_collection(QDRANT_COLLECTION)
    except Exception:
        qdrant_client.create_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=models.VectorParams(size=1024, distance=models.Distance.COSINE)
        )

# =========================
# ARCHIVE
# =========================
def archive_file(file_path: str, category: str):
    try:
        dest = os.path.join(PROCESSED_DIR, category)
        os.makedirs(dest, exist_ok=True)
        fname = os.path.basename(file_path)
        name, ext = os.path.splitext(fname)
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        new_name = f"{name}__{ts}{ext}"
        shutil.move(file_path, os.path.join(dest, new_name))
    except Exception as e:
        print(f"   ⚠️ Archive error: {e}")

# =========================
# MAIN ENGINE
# =========================
def process_single_file(file_path: str, doc_type: str):
    filename = os.path.basename(file_path)
    start_time = time.time()
    print(f"   ⚙️ Engine Start: {filename}")

    # Start log
    try:
        log_id = pg_start_log(filename, doc_type)
    except Exception as e:
        print(f"❌ DB Start Error: {e}")
        return

    try:
        # Estrai chunks (page-aware per PDF)
        chunks = extract_file_chunks(file_path, log_id)
        chunks = [c for c in chunks if c.get("text_sem") and len(c["text_sem"]) >= MIN_CHUNK_LEN]

        if not chunks:
            elapsed = int((time.time() - start_time) * 1000)
            pg_close_log(log_id, "SKIPPED_EMPTY", 0, elapsed)
            return

        print(f"   🚀 Estratti {len(chunks)} chunk (doc_type={doc_type}).")
        print(f"   🧠 Embeddings batch_size={EMBED_BATCH_SIZE}, flush_size={DB_FLUSH_SIZE}, max_KG={MAX_KG_CHUNKS_PER_DOC}")

        # Buffers
        qdrant_points: List[models.PointStruct] = []
        pg_rows: List[Tuple] = []
        neo4j_jobs: List[Tuple] = []  # (chunk_uuid, chunk_index, toon_type, page_no, graph)

        # KG budget
        kg_done = 0

        # Process in embedding batches
        total_saved = 0
        chunk_index_global = 0

        for batch_start in range(0, len(chunks), EMBED_BATCH_SIZE):
            batch = chunks[batch_start: batch_start + EMBED_BATCH_SIZE]
            texts_sem = [b["text_sem"] for b in batch]

            # 1) Embeddings in batch
            vectors = embedder.encode(texts_sem, batch_size=EMBED_BATCH_SIZE, normalize_embeddings=True)

            # 2) Per chunk: prepare Qdrant + PG + (optional) KG
            for local_i, ch in enumerate(batch):
                text_raw = ch["text_raw"]
                text_sem = ch["text_sem"]
                toon_type = ch.get("toon_type") or "text"
                page_no = ch.get("page_no")

                chunk_uuid = str(uuid.uuid4())
                vec = vectors[local_i].tolist()

                # Qdrant
                payload = {
                    "filename": filename,
                    "type": doc_type,
                    "chunk_index": chunk_index_global,
                    "toon_type": toon_type,
                    "page_no": page_no,
                    "log_id": log_id
                }
                qdrant_points.append(
                    models.PointStruct(
                        id=chunk_uuid,
                        vector=vec,
                        payload=payload
                    )
                )

                # Postgres row
                meta = {
                    "qdrant_uuid": chunk_uuid,
                    "original_path": file_path,
                    "filename": filename,
                    "doc_type": doc_type,
                    "chunk_index": chunk_index_global,
                    "page_no": page_no,
                    "toon_type": toon_type
                }
                facts = extract_facts(text_sem)
                if facts:
                    meta["facts"] = facts
                pg_rows.append((log_id, chunk_index_global, toon_type, text_raw, text_sem, Json(meta)))

                # KG extraction: solo se candidato + budget disponibile
                if kg_done < MAX_KG_CHUNKS_PER_DOC and is_candidate_for_kg(text_sem, doc_type):
                    graph = extract_graph_data(text_sem, doc_type, filename)
                    if graph:
                        neo4j_jobs.append((chunk_uuid, chunk_index_global, toon_type, page_no, graph))
                        kg_done += 1

                total_saved += 1
                chunk_index_global += 1
                print(".", end="", flush=True)

                # Flush if needed
                if len(qdrant_points) >= DB_FLUSH_SIZE:
                    # Qdrant
                    qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_points)
                    # Postgres
                    flush_postgres_chunks_batch(pg_rows)
                    # Neo4j batch flush (UNWIND)
                    flush_neo4j_jobs_batch(neo4j_jobs, log_id, doc_type, filename)

                    qdrant_points.clear()
                    pg_rows.clear()
                    neo4j_jobs.clear()
                    print(" [Flush] ", end="", flush=True)

        # Final flush
        if qdrant_points:
            qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_points)
            flush_postgres_chunks_batch(pg_rows)
            flush_neo4j_jobs_batch(neo4j_jobs, log_id, doc_type, filename)
            print(" [Flush] ", end="", flush=True)

        elapsed = int((time.time() - start_time) * 1000)
        pg_close_log(log_id, "COMPLETED", total_saved, elapsed)
        print(f"\n✅ Completed: {filename} | chunks={total_saved} | KG_calls={kg_done} | ms={elapsed}")

        archive_file(file_path, doc_type)

    except Exception as e:
        elapsed = int((time.time() - start_time) * 1000)
        print(f"\n❌ Critical Error: {e}")
        pg_close_log(log_id, "FAILED", 0, elapsed, str(e))


def ensure_qdrant_collection():
    dim = embedder.get_sentence_embedding_dimension()

    try:
        info = qdrant_client.get_collection(QDRANT_COLLECTION)
        existing_dim = info.config.params.vectors.size

        if existing_dim != dim:
            raise RuntimeError(
                f"Qdrant collection '{QDRANT_COLLECTION}' dimension mismatch: "
                f"collection={existing_dim}, embedder={dim}. "
                f"Delete/recreate the collection or change name."
            )

        print(f"✅ Qdrant collection '{QDRANT_COLLECTION}' ready (dim={dim})")
        return

    except Exception:
        print(f"🆕 Creating Qdrant collection '{QDRANT_COLLECTION}' (dim={dim})")
        qdrant_client.create_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=models.VectorParams(
                size=dim,
                distance=models.Distance.COSINE
            )
        )


def main():
    os.makedirs(INBOX_DIR, exist_ok=True)
    os.makedirs(PROCESSED_DIR, exist_ok=True)

    ensure_qdrant_collection()

    print("=== Ingestion Engine START ===")
    print(f"INBOX: {INBOX_DIR}")

    # scan inbox
    for root, _, files in os.walk(INBOX_DIR):
        for fname in files:
            ext = os.path.splitext(fname)[1].lower().lstrip(".")
            category = None
            for cat, exts in CATEGORIES.items():
                if ext in exts:
                    category = cat
                    break
            if not category:
                continue

            fpath = os.path.join(root, fname)
            process_single_file(fpath, category)

    print("=== Ingestion Engine END ===")

    try:
        neo4j_driver.close()
    except Exception:
        pass

    try:
        pg_pool.closeall()
    except Exception:
        pass

if __name__ == "__main__":
    main()
