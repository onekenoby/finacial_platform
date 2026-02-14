import os
import json
import uuid
import glob
import shutil
import time
import datetime
import hashlib
import base64
import re
import psycopg2
from psycopg2.extras import Json, execute_values
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from typing import List, Dict, Tuple
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient, models
from neo4j import GraphDatabase
from openai import OpenAI

# --- CONFIGURAZIONE ---
BASE_DATA_DIR = "./data_ingestion"
INBOX_DIR = "./data_ingestion/INBOX"
PROCESSED_DIR = "./data_ingestion/processed"
BATCH_SIZE = 10 

# DB Configs
PG_DSN = "dbname=ai_ingestion user=admin password=admin_password host=localhost"
QDRANT_HOST = "localhost"
QDRANT_COLLECTION = "financial_docs"
NEO4J_URI = "bolt://localhost:7687"
NEO4J_AUTH = ("neo4j", "password_sicura")

# AI Configs
LM_STUDIO_URL = "http://localhost:1234/v1"
LM_STUDIO_API_KEY = "lm-studio"
MODEL_LLM_ID = "gemma-3-12b"
MODEL_VISION_ID = "qwen3-vl-8b"
MODEL_EMBEDDING_ID = "BAAI/bge-m3"

# --- 1. TASSONOMIE COMPLETE (11 CATEGORIE) ---
ONTOLOGIES = {
    "financial": {
        "role": "Financial Analyst",
        "nodes": '"Organization", "Metric", "Currency", "TimePeriod", "Asset"',
        "edges": '"REPORTED_VALUE", "INCREASED_BY", "DECREASED_BY", "PROJECTS", "OWNS"',
        "special_rule": "Focus on quantitative data, balance sheets, and P&L. Translate metrics to English (e.g., 'Utile' -> 'Profit')."
    },
    "client": {
        "role": "Wealth Manager",
        "nodes": '"Client", "Goal", "RiskProfile", "Asset", "FamilyMember"',
        "edges": '"HAS_GOAL", "HAS_RISK_PROFILE", "OWNS", "INTERESTED_IN", "RELATED_TO"',
        "special_rule": "Focus on KYC, personal goals, and preferences."
    },
    "risk": { # Risk & Compliance
        "role": "Chief Risk Officer",
        "nodes": '"Risk", "Regulation", "Control", "Audit", "Authority"',
        "edges": '"MITIGATES", "VIOLATES", "COMPLIES_WITH", "REGULATES", "CAUSES_RISK"',
        "special_rule": "Identify risks (Market, Credit, Operational) and regulatory frameworks."
    },
    "operations": {
        "role": "COO / Operations Manager",
        "nodes": '"Process", "Facility", "SupplyChain", "Efficiency", "Logistics"',
        "edges": '"OPTIMIZES", "REQUIRES", "LOCATED_AT", "DELIVERS", "MANAGES"',
        "special_rule": "Map operational workflows, logistics, and production."
    },
    "technology": {
        "role": "CTO / Tech Analyst",
        "nodes": '"Technology", "Software", "Hardware", "Protocol", "Innovation"',
        "edges": '"RUNS_ON", "UPGRADES", "DEPRECATES", "ENABLES", "UTILIZES"',
        "special_rule": "Focus on IT infrastructure, digital transformation, and software."
    },
    "strategy": { # Strategy & Management
        "role": "Chief Strategy Officer",
        "nodes": '"Strategy", "Vision", "Competitor", "Market", "Merger"',
        "edges": '"COMPETES_WITH", "TARGETS", "ACQUIRED", "PLANS_TO", "PARTNERS_WITH"',
        "special_rule": "Extract forward-looking statements, M&A, and competitive landscape."
    },
    "legal": { # Legal & Regulatory
        "role": "Legal Counsel",
        "nodes": '"Contract", "Law", "Clause", "Party", "Litigation"',
        "edges": '"SIGNS", "SUES", "ENFORCES", "PROHIBITS", "AGREES_TO"',
        "special_rule": "Focus on contractual obligations, laws, and legal disputes."
    },
    "products": { # Products & Services
        "role": "Product Manager",
        "nodes": '"Product", "Service", "Feature", "Pricing", "CustomerSegment"',
        "edges": '"COSTS", "INCLUDES", "LAUNCHED_IN", "SERVES", "REPLACES"',
        "special_rule": "Catalog products, services, and their attributes."
    },
    "educational": {
        "role": "Academic / Economist",
        "nodes": '"Concept", "Theory", "Indicator", "EconomicEvent", "Definition"',
        "edges": '"CAUSES", "DEFINED_AS", "CORRELATED_WITH", "EXPLAINS", "CONTRADICTS"',
        "special_rule": "Extract definitions, causal relationships, and theoretical concepts."
    },
    "sustainability": {
        "role": "ESG Analyst",
        "nodes": '"ESG_Factor", "Emission", "Initiative", "Standard", "Impact"',
        "edges": '"REDUCES", "IMPACTS", "ALIGNED_WITH", "TARGETS_NET_ZERO", "POLLUTES"',
        "special_rule": "Focus on Environmental, Social, and Governance (ESG) topics."
    },
    "generic": { # Fallback
        "role": "General Knowledge Analyst",
        "nodes": '"Entity", "Person", "Organization", "Location", "Concept", "Date"',
        "edges": '"RELATED_TO", "INVOLVES", "LOCATED_IN", "HAS_ATTRIBUTE", "MENTIONS"',
        "special_rule": "Extract general entities. Be broad and inclusive."
    }
}

# --- PROMPT TEMPLATES ---

MASTER_PROMPT_TEMPLATE = """
ROLE: {role}
TASK: Analyze the provided text chunk and extract a Knowledge Graph in JSON format.

--- CONTEXT INFORMATION ---
SOURCE FILENAME: "{filename}"
PROCESSING DATE: {today}
TEMPORAL HINT: {temporal_hint}

--- RIGID ONTOLOGY (Use ONLY these types in English) ---
NODES: {nodes}
CATCH-ALL NODE: "KeyTerm" (Use for critical concepts that don't fit above).
RELATIONS: {edges}
CATCH-ALL RELATION: "RELATED_TO".

--- RULES ---
1. {special_rule}
2. Output must be strictly VALID JSON.
3. Normalize all entity types and relation names to ENGLISH.
4. Keep specific proper names in original language.

--- EXAMPLE OUTPUT ---
{{"nodes": [{{"id": "EntityName", "type": "NodeType"}}], "edges": [{{"source": "EntityName", "target": "OtherEntity", "relation": "RELATION_TYPE"}}]}}
"""

CLASSIFIER_PROMPT = """
You are a Senior Document Classifier. 
Analyze the provided text preview and classify the document into ONE of these categories:

1. FINANCIAL (Balance sheets, P&L, annual reports, tax)
2. CLIENT (KYC, personal profiles, wealth goals)
3. RISK (Risk assessments, compliance, audit, anti-money laundering)
4. OPERATIONS (Logistics, supply chain, manufacturing, facilities)
5. TECHNOLOGY (IT specs, software manuals, digital transformation)
6. STRATEGY (Strategic plans, M&A, market analysis, competitor review)
7. LEGAL (Contracts, court rulings, regulations, terms of service)
8. PRODUCTS (Product sheets, service descriptions, pricing)
9. EDUCATIONAL (Textbooks, theories, research papers, training)
10. SUSTAINABILITY (ESG reports, carbon footprint, CSR)
11. GENERIC (General news, emails, unknown types)

Reply ONLY with the single category keyword (e.g., "RISK", "LEGAL", "GENERIC").
"""

# --- INIZIALIZZAZIONE ---
print("🔌 Connessione servizi...")
embedder = SentenceTransformer(MODEL_EMBEDDING_ID, device='cuda')
qdrant_client = QdrantClient(host=QDRANT_HOST, port=6333)
neo4j_driver = GraphDatabase.driver(NEO4J_URI, auth=NEO4J_AUTH)
llm_client = OpenAI(base_url=LM_STUDIO_URL, api_key=LM_STUDIO_API_KEY)

# --- UTILS ---
def calculate_hash(data_bytes):
    return hashlib.md5(data_bytes).hexdigest()

def build_dynamic_prompt(doc_type: str, filename: str) -> str:
    ontology = ONTOLOGIES.get(doc_type, ONTOLOGIES["generic"])
    match = re.search(r"(19|20)\d{2}", filename)
    temporal_hint = f"The filename contains '{match.group(0)}'. Assume unstated dates refer to this year." if match else "No specific year in filename."
    
    return MASTER_PROMPT_TEMPLATE.format(
        role=ontology["role"],
        filename=filename,
        today=datetime.datetime.now().strftime("%Y-%m-%d"),
        temporal_hint=temporal_hint,
        nodes=ontology["nodes"],
        edges=ontology["edges"],
        special_rule=ontology["special_rule"]
    )

# --- DISPATCHER LOGIC (Classificazione) ---

def get_preview_text(file_path: str) -> str:
    try:
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        if ext == ".pdf":
            doc = fitz.open(file_path)
            for i in range(min(2, len(doc))): text += doc[i].get_text()
            doc.close()
        elif ext in [".docx", ".doc"]:
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs[:20]])
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                if i > 2: break
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f: text = f.read(3000)
        return text[:3000]
    except Exception as e:
        print(f"⚠️ Errore lettura preview: {e}")
        return ""

def classify_document(file_path: str) -> str:
    filename = os.path.basename(file_path)
    preview = get_preview_text(file_path)
    
    if not preview.strip():
        return "generic"

    try:
        response = llm_client.chat.completions.create(
            model=MODEL_LLM_ID,
            messages=[
                {"role": "system", "content": CLASSIFIER_PROMPT},
                {"role": "user", "content": f"FILENAME: {filename}\n\nTEXT PREVIEW:\n{preview}"}
            ],
            temperature=0.0,
            max_tokens=15
        )
        cat = response.choices[0].message.content.strip().upper()
        
        # Mapping rigoroso
        if "RISK" in cat or "COMPLIANCE" in cat: return "risk"
        if "STRATEGY" in cat: return "strategy"
        if "LEGAL" in cat or "REGULA" in cat: return "legal"
        if "PRODUCT" in cat or "SERVICE" in cat: return "products"
        if "SUSTAINABILITY" in cat or "ESG" in cat: return "sustainability"
        if "OPERATION" in cat: return "operations"
        if "TECHNOLOGY" in cat: return "technology"
        if "EDUCATION" in cat: return "educational"
        if "CLIENT" in cat: return "client"
        if "FINANCIAL" in cat: return "financial"
        
        return "generic"
    except Exception as e:
        print(f"⚠️ Errore classificazione AI: {e}")
        return "generic"

# --- BATCH DB HANDLERS (Con Safety Patch) ---

def flush_postgres_batch(batch_data: List[Tuple]):
    if not batch_data: return
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    try:
        execute_values(cur, """
            INSERT INTO document_chunks (log_id, chunk_index, toon_type, content_raw, content_semantic, metadata_json)
            VALUES %s
        """, batch_data)
        conn.commit()
    except Exception as e:
        conn.rollback()
        print(f"   ⚠️ Postgres Batch Error: {e}")
    finally:
        cur.close()
        conn.close()

def flush_neo4j_batch(batch_graph_data: List[Dict], doc_type: str):
    """Scrittura su Neo4j con controlli di sicurezza anti-crash"""
    if not batch_graph_data: return
    with neo4j_driver.session() as session:
        for item in batch_graph_data:
            chunk_id = item['chunk_id']
            
            # 1. Nodi (Safety Check)
            for node in item.get("nodes", []):
                # Controllo chiavi mancanti
                if not isinstance(node, dict) or "id" not in node or "type" not in node: 
                    continue 
                try:
                    session.run("""
                        MERGE (c:Chunk {id: $cid, type: $dtype})
                        MERGE (n:Entity {name: $name}) ON CREATE SET n.type = $type 
                        MERGE (n)-[:MENTIONED_IN]->(c)
                    """, cid=chunk_id, dtype=doc_type, name=node['id'], type=node['type'])
                except Exception as e:
                    pass # Ignora singolo nodo fallito

            # 2. Archi (Safety Check)
            for edge in item.get("edges", []):
                # Controllo chiavi mancanti
                if not isinstance(edge, dict) or "source" not in edge or "target" not in edge or "relation" not in edge: 
                    continue
                
                try:
                    # Sanitize relation name (solo caratteri validi)
                    raw_rel = edge['relation'].upper().replace(" ", "_")
                    rel_type = "".join(c for c in raw_rel if c.isalnum() or c == '_')
                    if not rel_type: rel_type = "RELATED_TO"

                    session.run(f"""
                        MATCH (a:Entity {{name: $s}}), (b:Entity {{name: $t}})
                        MERGE (a)-[:{rel_type}]->(b)
                    """, s=edge['source'], t=edge['target'])
                except Exception as e:
                    pass # Ignora singolo arco fallito

# --- SINGLE DB HANDLERS & AI CORE ---

def pg_start_log(filename: str, source_type: str) -> int:
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    try:
        cur.execute("INSERT INTO ingestion_logs (source_name, source_type, status, ingestion_ts) VALUES (%s, %s, 'PROCESSING', NOW()) RETURNING log_id", (filename, source_type))
        log_id = cur.fetchone()[0]
        conn.commit()
        return log_id
    except Exception as e: conn.rollback(); raise e
    finally: cur.close(); conn.close()

def pg_save_image(log_id: int, image_bytes: bytes, mime_type: str, description: str) -> int:
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    img_hash = calculate_hash(image_bytes)
    try:
        cur.execute("SELECT image_id FROM ingestion_images WHERE image_hash = %s", (img_hash,))
        res = cur.fetchone()
        if res: return res[0]
        cur.execute("INSERT INTO ingestion_images (log_id, image_data, image_hash, mime_type, description_ai, ingestion_ts) VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING image_id", (log_id, psycopg2.Binary(image_bytes), img_hash, mime_type, description))
        image_id = cur.fetchone()[0]
        conn.commit()
        return image_id
    except: return -1
    finally: cur.close(); conn.close()

def pg_close_log(log_id: int, status: str, total_chunks: int, processing_ms: int, error_msg: str = None):
    conn = psycopg2.connect(PG_DSN)
    cur = conn.cursor()
    try:
        cur.execute("UPDATE ingestion_logs SET status = %s, total_chunks = %s, processing_time_ms = %s, error_message = %s WHERE log_id = %s", (status, total_chunks, processing_ms, error_msg, log_id))
        conn.commit()
    except: pass
    finally: cur.close(); conn.close()

def analyze_image_with_vision(image_bytes: bytes) -> str:
    try:
        base64_image = base64.b64encode(image_bytes).decode('utf-8')
        response = llm_client.chat.completions.create(
            model=MODEL_VISION_ID,
            messages=[{"role": "system", "content": "You are a financial analyst."},
                      {"role": "user", "content": [{"type": "text", "text": "Describe this image in ENGLISH."}, {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}]}],
            temperature=0.1, max_tokens=300)
        return response.choices[0].message.content
    except: return "Description not available."

def extract_graph_data(text: str, doc_type: str, filename: str) -> Dict:
    system_prompt = build_dynamic_prompt(doc_type, filename)
    try:
        response = llm_client.chat.completions.create(
            model=MODEL_LLM_ID,
            messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": f"TEXT:\n{text}"}], temperature=0.0)
        content = response.choices[0].message.content.strip()
        # Pulizia robusta del Markdown JSON
        if "```json" in content: content = content.split("```json")[1].split("```")[0].strip()
        elif "```" in content: content = content.split("```")[1].split("```")[0].strip()
        return json.loads(content)
    except: return None

# --- FILE PARSERS ---
def extract_pdf_content(file_path: str, log_id: int) -> str:
    full_text = ""
    doc = fitz.open(file_path)
    print(f"   🔍 Analisi PDF ({len(doc)} pagine)...")
    for page in doc:
        full_text += page.get_text() + "\n"
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            if len(base_image["image"]) < 3000: continue
            print("   👁️ Vision AI processing...")
            desc = analyze_image_with_vision(base_image["image"])
            img_id = pg_save_image(log_id, base_image["image"], f"image/{base_image['ext']}", desc)
            full_text += f"\n\n[[IMAGE_DESC_AI (ID: {img_id}): {desc}]]\n\n"
    doc.close()
    return full_text

def get_file_content(file_path: str, log_id: int) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf": return extract_pdf_content(file_path, log_id)
    elif ext in [".docx", ".doc"]: return "\n".join([p.text for p in Document(file_path).paragraphs])
    elif ext in [".pptx", ".ppt"]: 
        text = []
        for slide in Presentation(file_path).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text: text.append(shape.text)
        return "\n".join(text)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f: return f.read()
    return ""

# --- MAIN PROCESSOR (Engine) ---

def process_single_file(file_path: str, doc_type: str):
    filename = os.path.basename(file_path)
    start_time = time.time()
    print(f"   ⚙️ Engine Start: {filename}")
    
    try: log_id = pg_start_log(filename, doc_type)
    except Exception as e: print(f"❌ DB Start Error: {e}"); return

    try:
        content = get_file_content(file_path, log_id)
        if not content.strip(): pg_close_log(log_id, "SKIPPED_EMPTY", 0, int((time.time()-start_time)*1000)); return

        chunk_size = 1000
        chunks_text = [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]
        qdrant_buffer, postgres_buffer, neo4j_buffer = [], [], []
        count = 0
        
        print(f"   🚀 Processing {len(chunks_text)} chunks (Batch Size: {BATCH_SIZE})...")

        for i, text_chunk in enumerate(chunks_text):
            if len(text_chunk) < 50: continue
            chunk_uuid = str(uuid.uuid4())
            toon_type = "image_description" if "[[IMAGE_DESC_AI" in text_chunk else "text"
            
            # Embedding
            vec = embedder.encode(text_chunk, normalize_embeddings=True).tolist()
            qdrant_buffer.append(models.PointStruct(id=chunk_uuid, vector=vec, payload={"filename": filename, "type": doc_type, "chunk_index": i, "toon_type": toon_type}))
            
            # Graph
            g_data = extract_graph_data(text_chunk, doc_type, filename)
            if g_data: g_data['chunk_id'] = chunk_uuid; neo4j_buffer.append(g_data)
            
            # Postgres
            meta = {"qdrant_uuid": chunk_uuid, "original_path": file_path}
            postgres_buffer.append((log_id, i, toon_type, text_chunk, text_chunk, Json(meta)))
            
            count += 1
            print(f".", end="", flush=True)

            if len(qdrant_buffer) >= BATCH_SIZE:
                qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_buffer)
                flush_postgres_batch(postgres_buffer)
                flush_neo4j_batch(neo4j_buffer, doc_type)
                qdrant_buffer, postgres_buffer, neo4j_buffer = [], [], []
                print(f" [Flush] ", end="")

        if qdrant_buffer:
            qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_buffer)
            flush_postgres_batch(postgres_buffer)
            flush_neo4j_batch(neo4j_buffer, doc_type)
            print(" [Flush] ", end="")

        elapsed = int((time.time() - start_time) * 1000)
        pg_close_log(log_id, "COMPLETED", count, elapsed)
        print(f"\n✅ Completed: {filename}")
        archive_file(file_path, doc_type)

    except Exception as e:
        elapsed = int((time.time() - start_time) * 1000)
        print(f"\n❌ Critical Error: {e}")
        pg_close_log(log_id, "FAILED", 0, elapsed, str(e))

def archive_file(file_path: str, category: str):
    try:
        dest = os.path.join(PROCESSED_DIR, category)
        os.makedirs(dest, exist_ok=True)
        fname = os.path.basename(file_path)
        name, ext = os.path.splitext(fname)
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        shutil.move(file_path, os.path.join(dest, f"{name}_{ts}{ext}"))
        print(f"   📦 Archived to {category}")
    except Exception as e: print(f"   ❌ Archive Error: {e}")

def main():
    if not os.path.exists(INBOX_DIR): os.makedirs(INBOX_DIR)
    if not qdrant_client.collection_exists(QDRANT_COLLECTION):
        qdrant_client.create_collection(collection_name=QDRANT_COLLECTION, vectors_config=models.VectorParams(size=1024, distance=models.Distance.COSINE))

    files = glob.glob(os.path.join(INBOX_DIR, "*"))
    files = [f for f in files if os.path.isfile(f) and not f.endswith(".tmp")]
    
    if not files: print(f"📭 INBOX vuota ({INBOX_DIR})"); return
    print(f"📬 Trovati {len(files)} file in INBOX.")
    
    for f in files:
        print(f"\n🕵️ Classificazione: {os.path.basename(f)}...")
        doc_type = classify_document(f)
        print(f"   🏷️ Categoria rilevata: {doc_type.upper()}")
        process_single_file(f, doc_type)

    print("\n✅ Ingestion Finished.")
    neo4j_driver.close()

if __name__ == "__main__":
    os.makedirs(INBOX_DIR, exist_ok=True)
    os.makedirs(PROCESSED_DIR, exist_ok=True)
    main()