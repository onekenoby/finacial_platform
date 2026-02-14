"""
Ingestion Engine (Production-oriented) - FINAL DEFINITIVE (Max Info + Max Neo4j Graph)
- Supporto: PDF, DOCX, PPTX, TXT
- DB: PostgreSQL (TimescaleDB), Qdrant (TOON), Neo4j (Deep Graph + APOC)
- AI locale (LM Studio): LLM (KG + classifier) + Vision (immagini PDF)
- Embeddings: SentenceTransformers (batch GPU)

OBIETTIVI (DEFINITIVI):
1) Massimizzare recupero informazione dai file: testo + immagini/tabelle/grafici + struttura (Document/Page/Chunk).
2) Massimizzare densità/qualità grafo Neo4j: più nodi/relazioni + più proprietà + provenance, pulite e interrogabili.

NOTE:
- Neo4j è "best effort": errori Neo4j non bloccano l'ingest.
- Qdrant collection: non viene mai cancellata automaticamente (production-safe).
"""

import os
import re
import json
import time
import uuid
import shutil
import hashlib
import base64
from typing import List, Dict, Tuple, Optional, Any

import fitz  # PyMuPDF
import psycopg2
from psycopg2.extras import Json, execute_values
from psycopg2.pool import SimpleConnectionPool

from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient, models
from neo4j import GraphDatabase
from openai import OpenAI

# =========================
# CONFIGURAZIONE (PROFILE: BALANCED SPEED/QUALITY)
# =========================
BASE_DATA_DIR = "./data_ingestion"
INBOX_DIR = os.path.join(BASE_DATA_DIR, "INBOX")
PROCESSED_DIR = os.path.join(BASE_DATA_DIR, "PROCESSED")
FAILED_DIR = os.path.join(BASE_DATA_DIR, "FAILED")

CATEGORIES = {
    "pdf": ["pdf"],
    "docx": ["docx", "doc"],
    "pptx": ["pptx", "ppt"],
    "txt": ["txt"]
}

# Chunking
CHUNK_MAX_CHARS = 1400
CHUNK_OVERLAP_CHARS = 200
MIN_CHUNK_LEN = 40

# KG Context
CONTEXT_WINDOW_CHARS = 260
INCLUDE_CONTEXT_IN_KG = True

# Batching
DB_FLUSH_SIZE = 128      # Ridotto per vedere progressi più spesso
EMBED_BATCH_SIZE = 32    # Ridotto per risparmiare VRAM per l'LLM

# KG Extraction (SPEED TUNED)
KG_MIN_LEN = 300
MAX_KG_CHUNKS_PER_DOC = 8  # Focus sui concetti chiave, non su tutto
KG_KEYWORDS = [
    "risk", "guidance", "forecast", "revenue", "earnings", "inflation", "rate", "spread",
    "debt", "cash", "profit", "loss", "agreement", "law", "covenant", "margin", "ebitda",
    "capex", "opex", "dividend", "default", "rating"
]

# Classifier (DISABLED FOR SPEED)
USE_LLM_CLASSIFIER = False # Usa solo Keyword matching (istantaneo)
CLASSIFIER_MAX_PER_DOC = 0

# Vision (SPEED TUNED)
PDF_VISION_ENABLED = True
PDF_VISION_ONLY_IF_TEXT_SCARSO = True
PDF_MIN_TEXT_LEN_FOR_NO_VISION = 400 # Meno aggressivo
PDF_MAX_IMAGES_PER_PAGE = 1          # Solo l'immagine principale
PDF_MIN_IMAGE_BYTES = 50_000         # Ignora immagini piccole/sporche
PDF_MIN_IMAGE_DIM = 300
VISION_MAX_IMAGE_BYTES = 2_500_000

# Qdrant
QDRANT_HOST = os.getenv("QDRANT_HOST", "127.0.0.1")
QDRANT_PORT = int(os.getenv("QDRANT_PORT", "6333"))
QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", "financial_docs")

# Postgres
PG_HOST = os.getenv("PG_HOST", "127.0.0.1")
PG_PORT = int(os.getenv("PG_PORT", "5432"))
PG_DB = os.getenv("PG_DB", "ai_ingestion")
PG_USER = os.getenv("PG_USER", "admin")
PG_PASS = os.getenv("PG_PASS", "admin_password")
PG_MIN_CONN = int(os.getenv("PG_MIN_CONN", "1"))
PG_MAX_CONN = int(os.getenv("PG_MAX_CONN", "5"))

# Neo4j
NEO4J_URI = os.getenv("NEO4J_URI", "bolt://127.0.0.1:7687")
NEO4J_USER = os.getenv("NEO4J_USER", "neo4j")
NEO4J_PASS = os.getenv("NEO4J_PASS", "password_sicura")
NEO4J_ENABLED = os.getenv("NEO4J_ENABLED", "1") == "1"

# LM Studio / OpenAI-compatible endpoint
LM_BASE_URL = os.getenv("LM_BASE_URL", "http://127.0.0.1:1234/v1")
LM_API_KEY = os.getenv("LM_API_KEY", "lm-studio")

LLM_MODEL_NAME = os.getenv("LLM_MODEL_NAME", "google/gemma-3-12b")
VISION_MODEL_NAME = os.getenv("VISION_MODEL_NAME", LLM_MODEL_NAME)

EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "BAAI/bge-m3")

# LLM reliability
LLM_MAX_TOKENS = 1500
LLM_TEMPERATURE = 0.1
LLM_RETRIES = 2

# RAG payload policy
QDRANT_STORE_TEXT = True
QDRANT_TEXT_MAX_CHARS = 2200  # abbastanza per answer grounding, senza payload enormi

# =========================
# CLIENTS INIT
# =========================
openai_client = OpenAI(base_url=LM_BASE_URL, api_key=LM_API_KEY)
embedder = SentenceTransformer(EMBEDDING_MODEL_NAME)
qdrant_client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)

neo4j_driver = None
if NEO4J_ENABLED:
    try:
        neo4j_driver = GraphDatabase.driver(NEO4J_URI, auth=(NEO4J_USER, NEO4J_PASS))
    except Exception as _e:
        print(f"⚠️ Neo4j disabled (driver init failed): {_e}")
        NEO4J_ENABLED = False

pg_pool = SimpleConnectionPool(
    PG_MIN_CONN, PG_MAX_CONN,
    host=PG_HOST, port=PG_PORT, dbname=PG_DB,
    user=PG_USER, password=PG_PASS
)

# =========================
# PROMPTS
# =========================
CLASSIFIER_PROMPT = """
You are a strict classifier.
Given a text chunk (and optional context), decide if it is useful to extract entities/relations for a financial knowledge graph.
Return ONLY valid JSON:
{"useful": true/false, "reason": "...", "signals": ["..."]}

Be conservative: if uncertain, return useful=false.
"""

KG_PROMPT = """
You are a Financial Knowledge Engineer.
Extract a HIGH-FIDELITY Knowledge Graph from the provided text.

You MUST extract:
1) NODES:
   - Mandatory: "id", "label", "type" (PascalCase: Company, FinancialMetric, Risk, Instrument, Rate, Index, Person, Contract, Regulation, Date, Period, Geography, Sector, Product, Event, KPI, Scenario, Forecast, Guidance, Covenant...)
   - Mandatory: "properties" (flat JSON) with as many factual attributes as possible from the text.

2) EDGES:
   - Mandatory: "source", "target", "relation" (UPPER_SNAKE_CASE)
   - Mandatory: "properties" (flat JSON) with factual attributes.

CRITICAL RULES:
- Properties MUST be flat. Values must be strings, numbers, or booleans.
- Add "confidence" (low|medium|high) where applicable.
- Include normalized attributes where possible: currency, unit, period, date, value, value_min/value_max, yoy/qoq, basis_points.
- Do NOT invent facts. Only extract what is explicitly supported.

Return ONLY JSON:
{"nodes":[...],"edges":[...]}
"""

KG_ENRICH_PROMPT = """
You are a Financial Knowledge Engineer.
You receive:
(A) TEXT (chunk + optional context)
(B) PRELIMINARY_GRAPH JSON

TASK:
Enrich the graph with as many additional factual properties as possible for BOTH nodes and edges, ONLY if explicitly supported by TEXT.
Add provenance to BOTH nodes and edges where possible:
- source_doc (filename)
- source_page (integer if known)
- source_chunk_index (integer if known)
- source_snippet (<= 180 chars from TEXT)

Return ONLY JSON with same schema:
{"nodes":[...],"edges":[...]}
"""

VISION_PROMPT = """
You are a strict financial analyst specialized in reading charts, tables, and diagrams extracted from PDF documents.

Your task is to analyze the provided image and extract ONLY factual, visible information.

CRITICAL RULES:
- Do NOT guess or infer information that is not clearly visible.
- Do NOT add interpretations, causes, or explanations.
- If something is unclear, unreadable, or ambiguous, list it under "unreadable_parts".
- Be conservative. When in doubt, omit.

Return ONLY valid JSON.
No markdown.
No comments.
No explanations.

JSON SCHEMA (must be respected exactly):
{
  "kind": "table|chart|diagram|photo|other",
  "key_points": [
    "short factual statements describing what is visible"
  ],
  "numbers": [
    {
      "label": "what the number refers to",
      "value": "numeric value as string",
      "unit": "%",
      "currency": "EUR|USD|GBP|null",
      "period": "date or period if visible"
    }
  ],
  "entities": [
    {
      "type": "Company|Index|Metric|Rate|Country|Other",
      "label": "exact visible name"
    }
  ],
  "unreadable_parts": [
    "description of unreadable or unclear elements"
  ]
}

If the image is decorative, low quality, or contains no extractable data:
- set kind = "other"
- leave all arrays empty.
"""


VISION_RECONCILE_PROMPT = """
You are a strict financial analyst specialized in charts/tables extracted from PDFs.

You receive:
(A) PAGE_TEXT (raw text extracted from the SAME page)
(B) VISION_JSON (a first-pass structured extraction from the image)

TASK:
- Keep ONLY information supported by PAGE_TEXT OR clearly visible chart/table semantics in VISION_JSON.
- Remove hallucinations, vague claims, or items not grounded.
- Normalize numbers and units when possible (%, bps, EUR/USD, dates/periods).
- If something is unclear, move it to unreadable_parts.
- Be conservative.

Return ONLY valid JSON with the SAME schema:
{
  "kind": "table|chart|diagram|photo|other",
  "key_points": ["..."],
  "numbers": [{"label":"", "value":"", "unit":"", "currency":"", "period":""}],
  "entities": [{"type":"", "label":""}],
  "unreadable_parts": ["..."]
}
"""


# =========================
# UTILS
# =========================
def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()

def sha256_file(path: str, chunk_size: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            b = f.read(chunk_size)
            if not b:
                break
            h.update(b)
    return h.hexdigest()

def normalize_ws(text: str) -> str:
    text = (text or "").replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def normalize_entity_id(s: str) -> str:
    s = (s or "").strip().lower()
    s = s.replace('"', "").replace("'", "")
    s = re.sub(r"\s+", " ", s)
    return s[:180]

def safe_json_extract(raw: str) -> Optional[Dict]:
    if not raw:
        return None
    try:
        s = raw.strip()
        s = re.sub(r"^```(?:json)?", "", s, flags=re.IGNORECASE).strip()
        s = re.sub(r"```$", "", s).strip()
        if not s.startswith("{"):
            m = re.search(r"\{.*\}", s, flags=re.DOTALL)
            if m:
                s = m.group(0).strip()
        return json.loads(s)
    except Exception:
        return None

def split_text_with_overlap(text: str, max_chars: int, overlap: int) -> List[str]:
    text = normalize_ws(text)
    if len(text) <= max_chars:
        return [text]
    out = []
    i = 0
    step = max(1, max_chars - overlap)
    while i < len(text):
        out.append(text[i:i + max_chars])
        i += step
    return out

def split_paragraphs(text: str, max_chars: int, overlap: int) -> List[str]:
    text = normalize_ws(text)
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    merged = []
    buf = ""
    for p in paras:
        if len(buf) + len(p) + 2 <= max_chars:
            buf = (buf + "\n\n" + p).strip()
        else:
            if buf:
                merged.append(buf)
            buf = p
    if buf:
        merged.append(buf)

    final_chunks = []
    for m in merged:
        final_chunks.extend(split_text_with_overlap(m, max_chars, overlap))
    return final_chunks

def extract_facts(text: str) -> Dict:
    if not text:
        return {}
    t = text[:20000]
    perc = re.findall(r"\b\d+(?:[\.,]\d+)?\s?%\b", t)
    bps = re.findall(r"\b\d+(?:[\.,]\d+)?\s?bps\b", t, flags=re.IGNORECASE)
    currency = re.findall(r"(?:€\s?\d[\d\.,]*|\$\s?\d[\d\.,]*|£\s?\d[\d\.,]*|\b\d[\d\.,]*\s?(?:EUR|USD|GBP)\b)", t)
    dates = re.findall(r"\b\d{4}-\d{2}-\d{2}\b|\b\d{1,2}/\d{1,2}/\d{2,4}\b", t)

    def _uniq_limit(xs, limit=30):
        out, seen = [], set()
        for x in xs:
            x = x.strip()
            if not x or x in seen:
                continue
            seen.add(x)
            out.append(x)
            if len(out) >= limit:
                break
        return out

    facts = {}
    if perc: facts["percentages"] = _uniq_limit(perc)
    if bps: facts["bps"] = _uniq_limit(bps)
    if currency: facts["amounts"] = _uniq_limit(currency)
    if dates: facts["dates"] = _uniq_limit(dates)
    return facts

def find_section_hint(page_text: str) -> str:
    if not page_text:
        return ""
    lines = [ln.strip() for ln in page_text.splitlines() if ln.strip()]
    for ln in lines[:25]:
        if 4 <= len(ln) <= 90:
            if ln.count(".") <= 1 and ln.count(",") <= 2 and not ln.endswith("."):
                if re.search(r"[A-Za-z]", ln):
                    return ln[:90]
    return ""

def add_context_windows(chunks: List[Dict]) -> List[Dict]:
    if not chunks:
        return chunks
    texts = [c.get("text_sem", "") for c in chunks]
    for i, c in enumerate(chunks):
        prev_txt = texts[i - 1] if i > 0 else ""
        next_txt = texts[i + 1] if i + 1 < len(texts) else ""
        c["context_prev"] = prev_txt[-CONTEXT_WINDOW_CHARS:] if prev_txt else ""
        c["context_next"] = next_txt[:CONTEXT_WINDOW_CHARS] if next_txt else ""
    return chunks

def is_keyword_candidate(text: str) -> bool:
    t = (text or "").lower()
    if len(t) < KG_MIN_LEN:
        return False
    return any(k in t for k in KG_KEYWORDS)

# =========================
# POSTGRES (POOL)
# =========================
def pg_get_conn():
    return pg_pool.getconn()

def pg_put_conn(conn):
    pg_pool.putconn(conn)

def pg_start_log(source_name: str, source_type: str) -> int:
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO ingestion_logs (source_name, source_type, ingestion_ts, status) "
                "VALUES (%s, %s, NOW(), %s) RETURNING log_id",
                (source_name, source_type, "RUNNING")
            )
            log_id = cur.fetchone()[0]
        conn.commit()
        return log_id
    finally:
        pg_put_conn(conn)

def pg_close_log(log_id: int, status: str, total_chunks: int, processing_ms: int, error_msg: str = None):
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE ingestion_logs SET status = %s, total_chunks = %s, processing_time_ms = %s, error_message = %s "
                "WHERE log_id = %s",
                (status, total_chunks, processing_ms, error_msg, log_id)
            )
        conn.commit()
    except Exception:
        conn.rollback()
    finally:
        pg_put_conn(conn)

def pg_get_image_by_hash(image_hash: str, cur=None) -> Optional[Tuple[int, str]]:
    if cur is not None:
        cur.execute("SELECT image_id, description_ai FROM ingestion_images WHERE image_hash = %s LIMIT 1", (image_hash,))
        return cur.fetchone()
    conn = pg_get_conn()
    try:
        with conn.cursor() as _cur:
            _cur.execute("SELECT image_id, description_ai FROM ingestion_images WHERE image_hash = %s LIMIT 1", (image_hash,))
            return _cur.fetchone()
    except Exception:
        return None
    finally:
        pg_put_conn(conn)

def pg_save_image(log_id: int, image_bytes: bytes, mime_type: str, description: str, cur=None) -> int:
    img_hash = sha256_hex(image_bytes)
    cached = pg_get_image_by_hash(img_hash, cur=cur)
    if cached:
        return cached[0]

    sql = (
        "INSERT INTO ingestion_images (log_id, image_data, image_hash, mime_type, description_ai, ingestion_ts) "
        "VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING image_id"
    )
    params = (log_id, psycopg2.Binary(image_bytes), img_hash, mime_type, description)

    if cur is not None:
        cur.execute(sql, params)
        return cur.fetchone()[0]

    conn = pg_get_conn()
    try:
        with conn.cursor() as _cur:
            _cur.execute(sql, params)
            image_id = _cur.fetchone()[0]
        conn.commit()
        return image_id
    except Exception:
        conn.rollback()
        return -1
    finally:
        pg_put_conn(conn)

def flush_postgres_chunks_batch(batch_data: List[Tuple]):
    if not batch_data:
        return
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            execute_values(
                cur,
                """
                INSERT INTO document_chunks (log_id, chunk_index, toon_type, content_raw, content_semantic, metadata_json)
                VALUES %s
                """,
                batch_data
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        print(f"   ⚠️ Postgres Batch Error: {e}")
    finally:
        pg_put_conn(conn)

# =========================
# QDRANT
# =========================
def ensure_qdrant_collection():
    dim = embedder.get_sentence_embedding_dimension()
    try:
        info = qdrant_client.get_collection(QDRANT_COLLECTION)
        existing_dim = info.config.params.vectors.size
        if existing_dim != dim:
            raise RuntimeError(
                f"❌ Qdrant collection '{QDRANT_COLLECTION}' dimension mismatch: "
                f"collection={existing_dim}, embedder={dim}. "
                f"Action: change QDRANT_COLLECTION or recreate manually."
            )
        print(f"✅ Qdrant collection '{QDRANT_COLLECTION}' ready (dim={dim})")
    except RuntimeError:
        raise
    except Exception:
        print(f"🆕 Creating Qdrant collection '{QDRANT_COLLECTION}' (dim={dim})")
        qdrant_client.create_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=models.VectorParams(size=dim, distance=models.Distance.COSINE)
        )

# =========================
# NEO4J (DEEP GRAPH + STRUCTURE)
# =========================
NEO4J_BATCH_QUERY = """
UNWIND $rows AS r

MERGE (d:Document {doc_id: r.doc_id})
SET d.filename = r.filename,
    d.doc_type = r.doc_type,
    d.log_id = r.log_id,
    d.doc_sha256 = r.doc_sha256,
    d.ingested_at = datetime()

WITH d, r
FOREACH (_ IN CASE WHEN r.page_no IS NULL THEN [] ELSE [1] END |
  MERGE (p:Page {doc_id: r.doc_id, page_no: r.page_no})
  MERGE (d)-[:HAS_PAGE]->(p)
)

WITH d, r
MERGE (c:Chunk {id: r.chunk_id})
SET c.filename = r.filename,
    c.doc_type = r.doc_type,
    c.log_id = r.log_id,
    c.chunk_index = r.chunk_index,
    c.toon_type = r.toon_type,
    c.page = r.page_no,
    c.section_hint = r.section_hint,
    c.doc_id = r.doc_id,
    c.doc_sha256 = r.doc_sha256,
    c.ingested_at = datetime()

MERGE (d)-[:HAS_CHUNK]->(c)

WITH r, c
FOREACH (_ IN CASE WHEN r.page_no IS NULL THEN [] ELSE [1] END |
  MERGE (p:Page {doc_id: r.doc_id, page_no: r.page_no})
  MERGE (p)-[:HAS_CHUNK]->(c)
)

WITH r, c
UNWIND r.nodes AS n
CALL apoc.merge.node([n.type, "Entity"], {id: n.id}, {label: n.label}, {}) YIELD node AS e
SET e += n.props
SET e.source_doc = coalesce(e.source_doc, r.filename),
    e.last_seen_at = datetime(),
    e.first_seen_at = coalesce(e.first_seen_at, datetime())
MERGE (e)-[:MENTIONED_IN]->(c)

WITH r
UNWIND r.edges AS rel
MATCH (s:Entity {id: rel.source})
MATCH (t:Entity {id: rel.target})
CALL apoc.merge.relationship(s, rel.relation, {}, {}, t, {}) YIELD rel AS r_out
SET r_out += rel.props
SET r_out.source_doc = coalesce(r_out.source_doc, r.filename),
    r_out.last_seen_at = datetime(),
    r_out.first_seen_at = coalesce(r_out.first_seen_at, datetime())

RETURN count(*)
"""

def _parse_numeric_value(s: str) -> Tuple[Optional[float], Optional[str]]:
    if not isinstance(s, str):
        return None, None
    x = s.strip()
    if not x:
        return None, None
    m = re.search(r"(-?\d+(?:[.,]\d+)?)\s*bps\b", x, flags=re.IGNORECASE)
    if m:
        return float(m.group(1).replace(",", ".")), "bps"
    m = re.search(r"(-?\d+(?:[.,]\d+)?)\s*%\b", x)
    if m:
        return float(m.group(1).replace(",", ".")), "percent"
    m = re.search(r"(-?\d+(?:[.,]\d+)?)(\s*)([KMB])\b", x, flags=re.IGNORECASE)
    if m:
        base = float(m.group(1).replace(",", "."))
        suf = m.group(3).upper()
        mult = {"K": 1e3, "M": 1e6, "B": 1e9}.get(suf, 1.0)
        return base * mult, f"scale_{suf}"
    m = re.search(r"-?\d+(?:[.,]\d+)?", x)
    if m:
        return float(m.group(0).replace(",", ".")), None
    return None, None

def _flat_props(props) -> Dict:
    if not isinstance(props, dict):
        return {}
    out = {}
    for k, v in props.items():
        if not isinstance(k, str):
            continue
        kk = re.sub(r"[^a-zA-Z0-9_]+", "_", k.strip()).strip("_")[:60]
        if not kk:
            continue
        if isinstance(v, (str, int, float, bool)):
            out[kk] = v
            if isinstance(v, str):
                num, unit = _parse_numeric_value(v)
                if num is not None:
                    out[f"{kk}_num"] = num
                if unit:
                    out[f"{kk}_unit_hint"] = unit
    return out

def _clean_type(t: str) -> str:
    raw_type = (t or "Entity").strip()
    clean_type = "".join(x.capitalize() for x in re.split(r"[_\s\-\.]+", raw_type) if x)
    ntype = re.sub(r"[^a-zA-Z0-9]", "", clean_type) or "Entity"
    if not re.match(r"^[A-Za-z]", ntype):
        ntype = "Entity"
    return ntype[:60]

def _clean_rel(r: str) -> str:
    rr = (r or "RELATED_TO").upper().strip()
    rr = re.sub(r"[^A-Z0-9_]+", "_", rr).strip("_")
    if not rr or not re.match(r"^[A-Z]", rr):
        rr = "RELATED_TO"
    return rr[:60]

def _sanitize_graph(graph: Dict) -> Tuple[List[Dict], List[Dict]]:
    if not isinstance(graph, dict):
        return [], []
    raw_nodes = graph.get("nodes") or []
    raw_edges = graph.get("edges") or []
    if not isinstance(raw_nodes, list) or not isinstance(raw_edges, list):
        return [], []

    nodes = []
    seen_ids = set()
    for n in raw_nodes:
        if not isinstance(n, dict):
            continue
        nid_raw = n.get("id") or n.get("label")
        nid = normalize_entity_id(nid_raw)
        if not nid or nid in seen_ids:
            continue
        ntype = _clean_type(n.get("type"))
        label = (n.get("label") or nid)[:200]
        props = _flat_props(n.get("properties", {}))
        nodes.append({"id": nid, "label": label, "type": ntype, "props": props})
        seen_ids.add(nid)

    edges = []
    seen_e = set()
    for e in raw_edges:
        if not isinstance(e, dict):
            continue
        src = normalize_entity_id(e.get("source"))
        tgt = normalize_entity_id(e.get("target"))
        if not src or not tgt:
            continue
        rel = _clean_rel(e.get("relation"))
        props = _flat_props(e.get("properties", {}))
        key = (src, tgt, rel)
        if key in seen_e:
            continue
        seen_e.add(key)
        edges.append({"source": src, "target": tgt, "relation": rel, "props": props})

    return nodes, edges

def flush_neo4j_rows_batch(rows: List[Dict]):
    if not NEO4J_ENABLED or not neo4j_driver:
        return
    if not rows:
        return
    try:
        with neo4j_driver.session() as session:
            session.run(NEO4J_BATCH_QUERY, rows=rows)
    except Exception as e:
        print(f"   ⚠️ Neo4j Batch Error (skipped): {e}")

# =========================
# LLM CALLS (RETRY)
# =========================
def llm_chat(prompt: str, user_text: str, model: str, max_tokens: int = LLM_MAX_TOKENS) -> str:
    last_err = None
    for _ in range(LLM_RETRIES + 1):
        try:
            resp = openai_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": prompt.strip()},
                    {"role": "user", "content": user_text.strip()}
                ],
                temperature=LLM_TEMPERATURE,
                max_tokens=max_tokens
            )
            return resp.choices[0].message.content or ""
        except Exception as e:
            last_err = e
            time.sleep(0.4)
    raise RuntimeError(f"LLM call failed after retries: {last_err}")

def llm_chat_multimodal(prompt: str, image_bytes: bytes, model: str, max_tokens: int = 700) -> str:
    """
    OpenAI-compatible multimodal call (text + image).
    Works if the LM Studio model supports image_url.
    """
    if not image_bytes:
        return ""
    b64_image = base64.b64encode(image_bytes).decode("utf-8")

    last_err = None
    for _ in range(LLM_RETRIES + 1):
        try:
            resp = openai_client.chat.completions.create(
                model=model,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt.strip()},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_image}"}}
                    ]
                }],
                temperature=0.1,
                max_tokens=max_tokens
            )
            return resp.choices[0].message.content or ""
        except Exception as e:
            last_err = e
            time.sleep(0.4)

    raise RuntimeError(f"Multimodal LLM call failed after retries: {last_err}")




def classify_chunk_useful(text: str) -> bool:
    out = llm_chat(CLASSIFIER_PROMPT, text, LLM_MODEL_NAME, max_tokens=300)
    js = safe_json_extract(out) or {}
    return bool(js.get("useful"))

def extract_graph_base(text_for_llm: str) -> Optional[Dict]:
    out = llm_chat(KG_PROMPT, text_for_llm, LLM_MODEL_NAME)
    return safe_json_extract(out)

def enrich_graph(text_for_llm: str, base_graph: Dict, filename: str, page_no: Optional[int], chunk_index: int) -> Optional[Dict]:
    payload = (
        "TEXT:\n"
        f"{text_for_llm}\n\n"
        "PRELIMINARY_GRAPH_JSON:\n"
        f"{json.dumps(base_graph, ensure_ascii=False)}\n\n"
        f"FILENAME: {filename}\n"
        f"PAGE_NO: {page_no}\n"
        f"CHUNK_INDEX: {chunk_index}\n"
    )
    out = llm_chat(KG_ENRICH_PROMPT, payload, LLM_MODEL_NAME)
    return safe_json_extract(out)

def analyze_image_with_vision(image_bytes: bytes) -> Optional[Dict]:
    if not image_bytes or len(image_bytes) > VISION_MAX_IMAGE_BYTES:
        return None

    strict_prompt = VISION_PROMPT.strip() + "\n\n" + (
        "IMPORTANT:\n"
        "- Return ONLY JSON (no markdown, no explanations).\n"
        "- If unsure, set kind='other' and keep arrays empty.\n"
    )

    try:
        out = llm_chat_multimodal(strict_prompt, image_bytes, VISION_MODEL_NAME, max_tokens=700)
        js = safe_json_extract(out)
        if isinstance(js, dict) and "kind" in js and "key_points" in js:
            js.setdefault("numbers", [])
            js.setdefault("entities", [])
            js.setdefault("unreadable_parts", [])
            return js
        return None
    except Exception as e:
        print(f"   ⚠️ Vision skipped (non-blocking): {e}")
        return None


def reconcile_vision_with_page_text(page_text: str, vision_json: Dict, model: str) -> Optional[Dict]:
    """
    Second pass: align vision output with page text, remove hallucinations,
    normalize numbers/periods/units. Conservative.
    """
    if not isinstance(vision_json, dict):
        return None

    payload = (
        "PAGE_TEXT:\n"
        f"{normalize_ws(page_text)[:8000]}\n\n"
        "VISION_JSON:\n"
        f"{json.dumps(vision_json, ensure_ascii=False)}\n"
    )

    out = llm_chat(VISION_RECONCILE_PROMPT, payload, model, max_tokens=700)
    js = safe_json_extract(out)
    if isinstance(js, dict) and "kind" in js and "key_points" in js:
        js.setdefault("numbers", [])
        js.setdefault("entities", [])
        js.setdefault("unreadable_parts", [])
        return js
    return None


# =========================
# EXTRACTION LOGIC
# =========================
def extract_docx_chunks(file_path: str) -> List[Dict]:
    try:
        import docx
    except ImportError:
        return []
    doc = docx.Document(file_path)
    paras = [normalize_ws(p.text) for p in doc.paragraphs if p.text and p.text.strip()]

    out = []
    current_hint = ""
    for p in paras:
        possible_hint = find_section_hint(p)
        if possible_hint:
            current_hint = possible_hint
        for c in split_text_with_overlap(p, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
            out.append({"text_raw": c, "text_sem": c, "page_no": None, "toon_type": "text", "section_hint": current_hint})
    return out

def extract_pptx_chunks(file_path: str) -> List[Dict]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    prs = Presentation(file_path)
    out = []
    for i, slide in enumerate(prs.slides):
        buf = []
        slide_title = ""
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = normalize_ws(slide.shapes.title.text)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                clean_t = shape.text.strip()
                if clean_t:
                    buf.append(clean_t)
                    if not slide_title:
                        hint = find_section_hint(clean_t)
                        if hint:
                            slide_title = hint
        s = normalize_ws("\n".join(buf))
        if s:
            for c in split_text_with_overlap(s, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
                out.append({"text_raw": c, "text_sem": c, "page_no": i + 1, "toon_type": "text", "section_hint": slide_title})
    return out

def extract_txt_chunks(file_path: str) -> List[Dict]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = normalize_ws(f.read())
    section = find_section_hint(text)
    return [{"text_raw": c, "text_sem": c, "page_no": None, "toon_type": "text", "section_hint": section}
            for c in split_text_with_overlap(text, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)]

def remove_repeated_lines(pages_text: List[str], min_occurrences: int = 3) -> List[str]:
    freq = {}
    for t in pages_text:
        for ln in set([ln.strip() for ln in t.splitlines() if ln.strip()]):
            freq[ln] = freq.get(ln, 0) + 1
    repeated = {ln for ln, c in freq.items() if c >= min_occurrences and len(ln) <= 120}
    cleaned = []
    for t in pages_text:
        lines = [ln for ln in t.splitlines() if ln.strip() not in repeated]
        cleaned.append("\n".join(lines))
    return cleaned


def extract_pdf_chunks(file_path: str, log_id: int) -> List[Dict]:
    """
    PDF ingestion (definitiva):
    - Vision SOLO su pagine candidate (keyword figure/chart/table ecc.) + fallback "testo scarso + immagini"
    - Estrae immagini anche quando page.get_images non basta (fallback page.get_image_info)
    - Downscale + compressione JPEG prima del Vision (migliora detection e rientra nei limiti bytes)
    - Log immagini in Postgres (ingestion_images) con cache su hash
    - Genera chunk SEPARATI:
        - testo pagina: toon_type="text"
        - descrizione immagine: toon_type="image_description"
      => così le descrizioni vanno in embeddings e possono essere processate per Neo4j come chunk normali
    """
    try:
        from PIL import Image
        import io
    except ImportError:
        raise RuntimeError("Missing dependency: Pillow. Install with: pip install pillow")

    doc = fitz.open(file_path)
    pg_conn = None
    pg_cur = None

    # ====== Candidate keywords (IT+EN) ======
    CANDIDATE_PAT = re.compile(
        r"\b("
        r"figura|fig\.|grafico|chart|graph|plot|diagramma|diagram|schema|flowchart|"
        r"tabella|table|tavola|matrice|matrix|istogramma|bar\s*chart|line\s*chart|"
        r"scatter|box\s*plot|heatmap"
        r")\b",
        flags=re.IGNORECASE
    )

    # ====== Soft thresholds (per non perdere grafici “leggeri”) ======
    MIN_BYTES_SOFT = 8_000
    MIN_DIM_SOFT = 140
    MAX_IMAGES_PER_PAGE = max(1, int(PDF_MAX_IMAGES_PER_PAGE))

    # ====== Vision gating ======
    # Se vuoi STRICT solo keyword, metti False
    ALLOW_FALLBACK_SCARCE_TEXT = True
    TEXT_LEN_FOR_NO_VISION = int(PDF_MIN_TEXT_LEN_FOR_NO_VISION)

    # ====== P5000-friendly Vision image preparation ======
    MAX_VISION_LONG_SIDE = 1200
    JPEG_QUALITY_START = 82
    MAX_SEND_BYTES = min(int(VISION_MAX_IMAGE_BYTES), 1_600_000)

    def _get_pg_cur():
        nonlocal pg_conn, pg_cur
        if pg_conn is None:
            pg_conn = pg_get_conn()
            pg_cur = pg_conn.cursor()
        return pg_cur

    def _image_is_big_enough(w: int, h: int, nbytes: int) -> bool:
        # grafici piccoli ma leggibili: soglie soft
        if w and h and (w < MIN_DIM_SOFT or h < MIN_DIM_SOFT):
            return False
        if nbytes and nbytes < MIN_BYTES_SOFT and (not w or not h or (w < 220 and h < 220)):
            return False
        return True

    def _downscale_and_compress_for_vision(img_bytes: bytes) -> Optional[bytes]:
        if not img_bytes:
            return None
        try:
            im = Image.open(io.BytesIO(img_bytes))

            # RGB
            if im.mode not in ("RGB", "L"):
                im = im.convert("RGB")
            elif im.mode == "L":
                im = im.convert("RGB")

            w, h = im.size
            long_side = max(w, h)

            # resize
            if long_side > MAX_VISION_LONG_SIDE:
                scale = MAX_VISION_LONG_SIDE / float(long_side)
                new_w = max(1, int(w * scale))
                new_h = max(1, int(h * scale))
                im = im.resize((new_w, new_h), resample=Image.LANCZOS)

            # compress
            q = JPEG_QUALITY_START
            out = io.BytesIO()
            im.save(out, format="JPEG", quality=q, optimize=True)
            data = out.getvalue()

            while len(data) > MAX_SEND_BYTES and q >= 50:
                q -= 8
                out = io.BytesIO()
                im.save(out, format="JPEG", quality=q, optimize=True)
                data = out.getvalue()

            if len(data) > MAX_SEND_BYTES:
                w, h = im.size
                im2 = im.resize((max(1, int(w * 0.78)), max(1, int(h * 0.78))), resample=Image.LANCZOS)
                out = io.BytesIO()
                im2.save(out, format="JPEG", quality=58, optimize=True)
                data = out.getvalue()

            if len(data) > MAX_SEND_BYTES:
                return None

            return data
        except Exception:
            return None

    def _extract_images_from_page(page: fitz.Page) -> List[Dict[str, Any]]:
        images = []
        seen_xref = set()

        # A) Standard
        try:
            imgs = page.get_images(full=True) or []
        except Exception:
            imgs = []

        for img in imgs:
            try:
                xref = img[0]
                if xref in seen_xref:
                    continue
                base = doc.extract_image(xref) or {}
                b = base.get("image")
                if not b:
                    continue
                seen_xref.add(xref)
                images.append({
                    "xref": xref,
                    "bytes": b,
                    "ext": (base.get("ext") or "png").lower(),
                    "width": int(base.get("width") or 0),
                    "height": int(base.get("height") or 0),
                })
            except Exception:
                continue

        # B) Fallback
        if not images:
            try:
                info_list = page.get_image_info(xrefs=True) or []
            except Exception:
                info_list = []
            for it in info_list:
                try:
                    xref = it.get("xref")
                    if not xref or xref in seen_xref:
                        continue
                    base = doc.extract_image(xref) or {}
                    b = base.get("image")
                    if not b:
                        continue
                    seen_xref.add(xref)
                    images.append({
                        "xref": xref,
                        "bytes": b,
                        "ext": (base.get("ext") or "png").lower(),
                        "width": int(it.get("width") or base.get("width") or 0),
                        "height": int(it.get("height") or base.get("height") or 0),
                    })
                except Exception:
                    continue

        return images

    def _is_candidate_page(page_text: str, has_images: bool) -> bool:
        txt = page_text or ""
        # Primary: keyword
        if CANDIDATE_PAT.search(txt):
            return True
        # Optional fallback: testo scarso + immagini
        if ALLOW_FALLBACK_SCARCE_TEXT and has_images and len(txt) < TEXT_LEN_FOR_NO_VISION:
            return True
        return False

    pages_text = remove_repeated_lines([normalize_ws(doc[i].get_text("text")) for i in range(len(doc))])
    out_chunks: List[Dict] = []
    
    # raccogli multimodal context per pagina
    page_multimodal_accumulator: Dict[int, List[Dict[str, Any]]] = {}

    print(f"   🔍 Analisi PDF ({len(doc)} pagine)...")

    for page_idx in range(len(doc)):
        page_no = page_idx + 1
        page = doc[page_idx]
        page_text = pages_text[page_idx]
        section_hint = find_section_hint(page_text)
        page_multimodal_accumulator[page_no] = []
        
        # 1) Sempre: chunk testo pagina
        for ch in split_paragraphs(normalize_ws(page_text), CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
            out_chunks.append({
                "text_raw": ch,
                "text_sem": ch,
                "page_no": page_no,
                "toon_type": "text",
                "section_hint": section_hint
            })

        # 2) Vision SOLO su pagine candidate
        try:
            quick_imgs = page.get_images(full=False) or []
            has_images = len(quick_imgs) > 0
        except Exception:
            has_images = False

        allow_vision = bool(PDF_VISION_ENABLED) and _is_candidate_page(page_text, has_images)
        if not allow_vision:
            continue

        imgs = _extract_images_from_page(page)

        # Selezione: le più grandi per area
        candidates = []
        for imd in imgs:
            b = imd["bytes"]
            w = int(imd.get("width") or 0)
            h = int(imd.get("height") or 0)

            # fallback dimension via Pixmap se manca
            if (w == 0 or h == 0) and imd.get("xref"):
                try:
                    pix = fitz.Pixmap(doc, imd["xref"])
                    w, h = pix.width, pix.height
                except Exception:
                    pass

            if not _image_is_big_enough(w, h, len(b)):
                continue

            area = (w or 1) * (h or 1)
            candidates.append((area, imd))

        candidates.sort(key=lambda x: x[0], reverse=True)
        candidates = candidates[:MAX_IMAGES_PER_PAGE]

#---------------------------------------------
        used = 0
        for _, imd in candidates:
            if used >= MAX_IMAGES_PER_PAGE:
                break
            try:
                img_bytes_orig = imd["bytes"]
                ext = imd["ext"]

                # Cache su hash ORIGINAL bytes
                img_hash = sha256_hex(img_bytes_orig)
                cached = pg_get_image_by_hash(img_hash, cur=_get_pg_cur())
                if cached:
                    img_id, desc_ai = cached[0], cached[1] or ""
                else:
                    # prepara per vision
                    img_bytes_for_vision = _downscale_and_compress_for_vision(img_bytes_orig)
                    if not img_bytes_for_vision:
                        continue

                vis = analyze_image_with_vision(img_bytes_for_vision)

                if vis and isinstance(vis, dict):
                    # SECOND PASS: reconcile with page text
                    vis2 = reconcile_vision_with_page_text(page_text, vis, VISION_MODEL_NAME)
                    vis_clean = vis2 if (vis2 and isinstance(vis2, dict)) else vis

                    kind = vis_clean.get("kind", "other")
                    key_points = vis_clean.get("key_points", []) or []
                    numbers = vis_clean.get("numbers", []) or []
                    entities = vis_clean.get("entities", []) or []
                    unreadable = vis_clean.get("unreadable_parts", []) or []

                    # descrizione compatta (per chunk + embeddings)
                    desc_ai = f"[{kind}] " + "; ".join([str(x) for x in key_points[:10]])

                    if numbers:
                        desc_ai += " | numbers: " + "; ".join([
                            f"{n.get('label','')}: {n.get('value','')}{n.get('unit','')}"
                            + (f" {n.get('currency','')}".strip() if n.get("currency") else "")
                            + (f" ({n.get('period','')})" if n.get("period") else "")
                            for n in numbers[:12]
                        ])

                    if entities:
                        desc_ai += " | entities: " + "; ".join([
                            f"{e.get('type','Entity')}:{e.get('label','')}" for e in entities[:12]
                        ])

                    if unreadable:
                        desc_ai += " | unreadable: " + "; ".join([str(x) for x in unreadable[:6]])

                    # salva per contesto multimodale pagina (per KG “gold chunk”)
                    page_multimodal_accumulator[page_no].append({
                        "image_kind": kind,
                        "key_points": key_points[:12],
                        "numbers": numbers[:18],
                        "entities": entities[:18],
                        "unreadable_parts": unreadable[:10],
                    })
                else:
                    desc_ai = ""

                    if not desc_ai.strip():
                        continue

                    # Log immagine ORIGINAL su PG
                    img_id = pg_save_image(
                        log_id,
                        img_bytes_orig,
                        f"image/{ext}",
                        desc_ai,
                        cur=_get_pg_cur()
                    )

                if img_id and img_id > 0 and desc_ai.strip():
                    # 3) Chunk SEPARATO: descrizione immagine (andrà in embeddings + KG)
                    text_desc = normalize_ws(f"IMAGE_DESC_AI (image_id={img_id}, page={page_no}): {desc_ai}")
                    for ch in split_text_with_overlap(text_desc, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
                        out_chunks.append({
                            "text_raw": ch,
                            "text_sem": ch,
                            "page_no": page_no,
                            "toon_type": "image_description",
                            "section_hint": section_hint
                        })
                    used += 1

            except Exception:
                continue
#----------------------------------------------
        if pg_conn is not None:
            try:
                pg_conn.commit()
            except Exception:
                pg_conn.rollback()

        # 4) Chunk "GOLD": testo pagina + insight immagini (multimodale)
        insights = page_multimodal_accumulator.get(page_no) or []
        if insights:
            # costruisce un contesto testuale solido e “KG-friendly”
            parts = []
            parts.append(f"PAGE_NO: {page_no}")
            if section_hint:
                parts.append(f"SECTION_HINT: {section_hint}")
            parts.append("PAGE_TEXT:")
            parts.append(normalize_ws(page_text)[:12000])

            parts.append("\nIMAGE_INSIGHTS (AI, reconciled):")
            for idx, it in enumerate(insights, start=1):
                parts.append(f"\n- IMAGE #{idx} kind={it.get('image_kind','other')}")
                kps = it.get("key_points") or []
                if kps:
                    parts.append("  key_points: " + "; ".join([str(x) for x in kps[:10]]))
                nums = it.get("numbers") or []
                if nums:
                    parts.append("  numbers: " + "; ".join([
                        f"{n.get('label','')}: {n.get('value','')}{n.get('unit','')}"
                        + (f" {n.get('currency','')}".strip() if n.get("currency") else "")
                        + (f" ({n.get('period','')})" if n.get("period") else "")
                        for n in nums[:12]
                    ]))
                ents = it.get("entities") or []
                if ents:
                    parts.append("  entities: " + "; ".join([
                        f"{e.get('type','Entity')}:{e.get('label','')}" for e in ents[:12]
                    ]))
                unr = it.get("unreadable_parts") or []
                if unr:
                    parts.append("  unreadable: " + "; ".join([str(x) for x in unr[:6]]))

            multimodal_text = normalize_ws("\n".join(parts))
            # chunk separato: va in embeddings + può essere scelto per KG extraction
            for ch2 in split_text_with_overlap(multimodal_text, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS):
                out_chunks.append({
                    "text_raw": ch2,
                    "text_sem": ch2,
                    "page_no": page_no,
                    "toon_type": "page_text_plus_images",
                    "section_hint": section_hint
                })


    doc.close()
    if pg_cur is not None:
        try:
            pg_cur.close()
        except Exception:
            pass
    if pg_conn is not None:
        pg_put_conn(pg_conn)

    return out_chunks


def extract_file_chunks(file_path: str, log_id: int) -> List[Dict]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".doc":
        raise RuntimeError("UNSUPPORTED_FORMAT_DOC: .doc legacy not supported without conversion (doc->docx).")

    if ext == ".pdf":
        chunks = extract_pdf_chunks(file_path, log_id)
    elif ext == ".docx":
        chunks = extract_docx_chunks(file_path)
    elif ext in [".pptx", ".ppt"]:
        chunks = extract_pptx_chunks(file_path)
    elif ext == ".txt":
        chunks = extract_txt_chunks(file_path)
    else:
        chunks = []
    return add_context_windows(chunks)

# =========================
# FILE ARCHIVE
# =========================
def archive_file(file_path: str, ok: bool = True):
    os.makedirs(PROCESSED_DIR, exist_ok=True)
    os.makedirs(FAILED_DIR, exist_ok=True)
    dest_dir = PROCESSED_DIR if ok else FAILED_DIR
    dest = os.path.join(dest_dir, os.path.basename(file_path))
    try:
        if os.path.exists(dest):
            base, ext = os.path.splitext(dest)
            dest = f"{base}_{int(time.time())}{ext}"
        shutil.move(file_path, dest)
    except Exception as e:
        print(f"⚠️ archive_file failed: {e}")

# =========================
# MAIN PROCESSING LOOP
# =========================
def process_single_file(file_path: str, doc_type: str):
    filename = os.path.basename(file_path)
    start_time = time.time()
    print(f"   ⚙️ Engine Start: {filename}")

    # ids
    try:
        doc_sha = sha256_file(file_path)
    except Exception:
        doc_sha = ""
    doc_id = f"{filename}::{doc_sha[:16]}" if doc_sha else filename

    # start log
    try:
        log_id = pg_start_log(filename, doc_type)
    except Exception as e:
        print(f"❌ DB Start Error: {e}")
        return

    total_saved = 0
    kg_done = 0
    classifier_done = 0

    try:
        chunks = extract_file_chunks(file_path, log_id)
        chunks = [c for c in chunks if c.get("text_sem") and len(c["text_sem"]) >= MIN_CHUNK_LEN]

        if not chunks:
            elapsed = int((time.time() - start_time) * 1000)
            pg_close_log(log_id, "SKIPPED_EMPTY", 0, elapsed)
            archive_file(file_path, ok=True)
            return

        print(f"   🚀 Estratti {len(chunks)} chunk (doc_type={doc_type}).")
        print(f"   🧠 Embeddings batch_size={EMBED_BATCH_SIZE}, flush_size={DB_FLUSH_SIZE}, max_KG={MAX_KG_CHUNKS_PER_DOC}")

        qdrant_points: List[models.PointStruct] = []
        pg_rows: List[Tuple] = []
        neo4j_rows: List[Dict] = []

        chunk_index_global = 0

        for batch_start in range(0, len(chunks), EMBED_BATCH_SIZE):
            batch = chunks[batch_start: batch_start + EMBED_BATCH_SIZE]
            texts_sem = [b["text_sem"] for b in batch]
            vectors = embedder.encode(texts_sem, batch_size=EMBED_BATCH_SIZE, normalize_embeddings=True)

            for local_i, ch in enumerate(batch):
                text_raw = ch["text_raw"]
                text_sem = ch["text_sem"]
                toon_type = ch.get("toon_type", "text")
                page_no = ch.get("page_no")
                section_hint = ch.get("section_hint", "")
                ctx_prev = ch.get("context_prev", "")
                ctx_next = ch.get("context_next", "")

                chunk_uuid = str(uuid.uuid4())
                vec = vectors[local_i].tolist()
                facts = extract_facts(text_sem)

                # KG text
                if INCLUDE_CONTEXT_IN_KG:
                    kg_text = (
                        f"FILENAME: {filename}\n"
                        f"DOC_TYPE: {doc_type}\n"
                        f"PAGE_NO: {page_no}\n"
                        f"SECTION_HINT: {section_hint}\n\n"
                        f"CONTEXT_PREV:\n{ctx_prev}\n\n"
                        f"CHUNK:\n{text_sem}\n\n"
                        f"CONTEXT_NEXT:\n{ctx_next}\n"
                    )
                else:
                    kg_text = text_sem

                # --- KG extraction ---
                graph_data = None
                should_try_kg = False
                
                
                if kg_done < MAX_KG_CHUNKS_PER_DOC and len(text_sem) >= KG_MIN_LEN:
                    # PRIORITÀ: chunk multimodale “gold”
                    if toon_type in ("page_text_plus_images",):
                        should_try_kg = True
                    elif is_keyword_candidate(text_sem):
                        should_try_kg = True
                    elif USE_LLM_CLASSIFIER and classifier_done < CLASSIFIER_MAX_PER_DOC:
                        classifier_done += 1
                        should_try_kg = classify_chunk_useful(kg_text)

                        

                toon_entities = []
                toon_relations = []

                if should_try_kg and kg_done < MAX_KG_CHUNKS_PER_DOC:
                    base_graph = extract_graph_base(kg_text)
                    if base_graph and isinstance(base_graph, dict):
                        enriched = enrich_graph(kg_text, base_graph, filename, page_no, chunk_index_global)
                        graph_data = enriched if (enriched and isinstance(enriched, dict)) else base_graph

                        nodes, edges = _sanitize_graph(graph_data)
                        if nodes or edges:
                            snippet = normalize_ws(text_sem)[:180]
                            for n in nodes:
                                n["props"].setdefault("source_doc", filename)
                                if page_no is not None:
                                    n["props"].setdefault("source_page", int(page_no))
                                n["props"].setdefault("source_chunk_index", int(chunk_index_global))
                                n["props"].setdefault("source_snippet", snippet)
                                n["props"].setdefault("canonical_name", n.get("label", ""))

                            for e in edges:
                                e["props"].setdefault("source_doc", filename)
                                if page_no is not None:
                                    e["props"].setdefault("source_page", int(page_no))
                                e["props"].setdefault("source_chunk_index", int(chunk_index_global))
                                e["props"].setdefault("evidence_snippet", snippet)

                            neo4j_rows.append({
                                "doc_id": doc_id,
                                "doc_sha256": doc_sha,
                                "filename": filename,
                                "doc_type": doc_type,
                                "log_id": log_id,
                                "chunk_id": chunk_uuid,
                                "chunk_index": int(chunk_index_global),
                                "toon_type": toon_type,
                                "page_no": int(page_no) if page_no else None,
                                "section_hint": section_hint[:120] if section_hint else "",
                                "nodes": nodes,
                                "edges": edges
                            })
                            kg_done += 1

                        # extraction for qdrant metadata
                        toon_entities = list({
                            (n.get("label") or n.get("id") or "")[:120]
                            for n in (graph_data.get("nodes") or [])
                            if isinstance(n, dict)
                        })
                        toon_entities = [x for x in toon_entities if x][:25]

                        toon_relations = list({
                            (e.get("relation") or "")[:80]
                            for e in (graph_data.get("edges") or [])
                            if isinstance(e, dict)
                        })
                        toon_relations = [x for x in toon_relations if x][:40]

                # --- QDRANT payload (RAG-ready) ---
                toon_payload = {
                    "qdrant_uuid": chunk_uuid,     # utile lato RAG
                    "filename": filename,
                    "doc_type": doc_type,
                    "doc_id": doc_id,
                    "doc_sha256": doc_sha,
                    "log_id": log_id,
                    "chunk_index": int(chunk_index_global),
                    "page_no": int(page_no) if page_no else None,
                    "toon_type": toon_type,
                    "section_hint": (section_hint or "")[:200],
                    "facts": facts,
                    "toon_entities": toon_entities,
                    "toon_relations": toon_relations,
                }

                if QDRANT_STORE_TEXT:
                    toon_payload["content_semantic"] = normalize_ws(text_sem)[:QDRANT_TEXT_MAX_CHARS]
                    toon_payload["content_hash"] = sha256_hex(normalize_ws(text_sem).encode("utf-8"))

                qdrant_points.append(models.PointStruct(id=chunk_uuid, vector=vec, payload=toon_payload))

                # --- POSTGRES ---
                meta = {
                    "qdrant_uuid": chunk_uuid,
                    "original_path": file_path,
                    "filename": filename,
                    "doc_type": doc_type,
                    "doc_id": doc_id,
                    "doc_sha256": doc_sha,
                    "chunk_index": int(chunk_index_global),
                    "page_no": int(page_no) if page_no else None,
                    "toon_type": toon_type,
                    "section_hint": section_hint,
                    "facts": facts,
                    "toon_entities": toon_entities,
                    "toon_relations": toon_relations,
                }
                pg_rows.append((log_id, int(chunk_index_global), toon_type, text_raw, text_sem, Json(meta)))

                total_saved += 1
                chunk_index_global += 1
                print(".", end="", flush=True)

                # flush
                if len(qdrant_points) >= DB_FLUSH_SIZE:
                    flush_postgres_chunks_batch(pg_rows)
                    qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_points)
                    flush_neo4j_rows_batch(neo4j_rows)

                    qdrant_points.clear()
                    pg_rows.clear()
                    neo4j_rows.clear()
                    print(" [Flush] ", end="", flush=True)

        # final flush
        if pg_rows:
            flush_postgres_chunks_batch(pg_rows)
        if qdrant_points:
            qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_points)
        if neo4j_rows:
            flush_neo4j_rows_batch(neo4j_rows)

        elapsed = int((time.time() - start_time) * 1000)
        pg_close_log(log_id, "COMPLETED", total_saved, elapsed)
        print(f"\n✅ Completed: {filename} | chunks={total_saved} | KG_calls={kg_done} | classifier_calls={classifier_done} | ms={elapsed}")
        archive_file(file_path, ok=True)

    except Exception as e:
        elapsed = int((time.time() - start_time) * 1000)
        msg = str(e)

        if msg.startswith("UNSUPPORTED_FORMAT_DOC"):
            print(f"\n⚠️ Skipped unsupported .doc: {filename}")
            pg_close_log(log_id, "SKIPPED_UNSUPPORTED_DOC", 0, elapsed, msg)
            archive_file(file_path, ok=False)
            return

        print(f"\n❌ Critical Error: {e}")
        pg_close_log(log_id, "FAILED", total_saved, elapsed, msg)
        archive_file(file_path, ok=False)

def main():
    os.makedirs(INBOX_DIR, exist_ok=True)
    os.makedirs(PROCESSED_DIR, exist_ok=True)
    os.makedirs(FAILED_DIR, exist_ok=True)

    ensure_qdrant_collection()

    print("=== Ingestion Engine START ===")
    print(f"INBOX: {INBOX_DIR}")

    for root, _, files in os.walk(INBOX_DIR):
        for fname in files:
            ext = os.path.splitext(fname)[1].lower().lstrip(".")
            category = None
            for cat, exts in CATEGORIES.items():
                if ext in exts:
                    category = cat
                    break
            if not category:
                continue
            process_single_file(os.path.join(root, fname), category)

    print("=== Ingestion Engine END ===")
    try:
        if neo4j_driver:
            neo4j_driver.close()
    except Exception:
        pass
    try:
        pg_pool.closeall()
    except Exception:
        pass

if __name__ == "__main__":
    main()
