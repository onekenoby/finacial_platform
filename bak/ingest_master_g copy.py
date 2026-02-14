"""
Ingestion Engine (Production-oriented)
- Supporto: PDF, DOCX, PPTX, TXT
- DB: PostgreSQL (TimescaleDB), Qdrant, Neo4j
- AI locale (LM Studio): LLM (KG + classifier) + Vision (immagini PDF)
- Embeddings locali: SentenceTransformers (batch su GPU)
- TOON: gestione chunk + metadata + (opzionale) semantic normalization

Migliorie applicate (rispetto al tuo codice originale):
1) Prompt più robusti (classifier + KG + vision) con vincoli anti-hallucination
2) Embeddings calcolati in batch (GPU sfruttata meglio)
3) Riduzione drastica chiamate LLM (KG solo su chunk candidati + limite max per documento)
4) Chunking migliore (PDF page-aware + split semantico con overlap; DOCX per paragrafi; PPTX per shape)
5) Neo4j write path ottimizzato (UNWIND batch; no query per nodo/arco)
   - Relationship uniforme :REL con property "type" per evitare rel-type dinamici (più sicuro e più batch-friendly)
6) Postgres connection pooling (meno overhead, più velocità e stabilità)
7) Vision caching/dedup su SHA-256 (evita rianalisi immagini ripetute)
8) content_semantic realmente "pulito" (normalizzazione whitespace + rimozione righe ripetute nei PDF)

NOTE:
- Assumo che la collection Qdrant "financial_docs" abbia vettori dimensione 1024 (BGE-M3).
- LM Studio endpoint OpenAI-compatible.
"""

import os
import re
import glob
import json
import uuid
import time
import shutil
import base64
import hashlib
import datetime
from typing import List, Dict, Tuple, Optional

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

import psycopg2
from psycopg2.extras import Json, execute_values
from psycopg2.pool import SimpleConnectionPool

from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient, models
from neo4j import GraphDatabase
from openai import OpenAI


# =========================
# CONFIGURAZIONE
# =========================
BASE_DATA_DIR = "./data_ingestion"
INBOX_DIR = os.path.join(BASE_DATA_DIR, "INBOX")
PROCESSED_DIR = os.path.join(BASE_DATA_DIR, "processed")

# Batch / performance
DB_FLUSH_SIZE = 50                 # flush (Qdrant + PG + Neo4j) ogni N chunk
EMBED_BATCH_SIZE = 32              # batch embeddings su GPU
MIN_CHUNK_LEN = 80                 # evita chunk troppo corti
MAX_KG_CHUNKS_PER_DOC = 60         # limita chiamate LLM per KG per documento

# Chunking (caratteri; per token servirebbe tokenizer specifico)
CHUNK_MAX_CHARS = 4000
CHUNK_OVERLAP_CHARS = 200

# PDF + Vision
PDF_MAX_PAGES_PREVIEW = 2
PDF_IMAGE_MIN_BYTES = 25_000       # filtro più robusto del solo >3000 bytes
PDF_IMAGE_MIN_DIM = 240            # width/height min (evita icone)
PDF_VISION_MAX_IMAGES_PER_PAGE = 3
PDF_VISION_ONLY_IF_TEXT_SCARSO = True
PDF_TEXT_SCARSO_THRESHOLD = 250    # se testo pagina < soglia, abilita Vision più facilmente

# DB Configs
PG_DSN = "dbname=ai_ingestion user=admin password=admin_password host=localhost"
PG_POOL_MINCONN = 1
PG_POOL_MAXCONN = 6

QDRANT_HOST = "localhost"
QDRANT_PORT = 6333
QDRANT_COLLECTION = "financial_docs"

NEO4J_URI = "bolt://localhost:7687"
NEO4J_AUTH = ("neo4j", "password_sicura")

# AI Configs (LM Studio OpenAI-compatible)
LM_STUDIO_URL = "http://localhost:1234/v1"
LM_STUDIO_API_KEY = "lm-studio"
MODEL_LLM_ID = "gemma-3-12b"
MODEL_VISION_ID = "qwen3-vl-8b"
MODEL_EMBEDDING_ID = "BAAI/bge-m3"

# Device embeddings: prova CUDA, fallback CPU
EMBED_DEVICE = "cuda"


# =========================
# ONTOLOGIES (11 categorie)
# =========================
ONTOLOGIES = {
    "financial": {
        "role": "Financial Analyst",
        "nodes": '"Organization", "Metric", "Currency", "TimePeriod", "Asset"',
        "edges": '"REPORTED_VALUE", "INCREASED_BY", "DECREASED_BY", "PROJECTS", "OWNS"',
        "special_rule": "Focus on quantitative data, balance sheets, and P&L. Translate metrics to English (e.g., 'Utile' -> 'Profit')."
    },
    "client": {
        "role": "Wealth Manager",
        "nodes": '"Client", "Goal", "RiskProfile", "Asset", "FamilyMember"',
        "edges": '"HAS_GOAL", "HAS_RISK_PROFILE", "OWNS", "INTERESTED_IN", "RELATED_TO"',
        "special_rule": "Focus on KYC, personal goals, and preferences."
    },
    "risk": {
        "role": "Chief Risk Officer",
        "nodes": '"Risk", "Regulation", "Control", "Audit", "Authority"',
        "edges": '"MITIGATES", "VIOLATES", "COMPLIES_WITH", "REGULATES", "CAUSES_RISK"',
        "special_rule": "Identify risks (Market, Credit, Operational) and regulatory frameworks."
    },
    "operations": {
        "role": "COO / Operations Manager",
        "nodes": '"Process", "Facility", "SupplyChain", "Efficiency", "Logistics"',
        "edges": '"OPTIMIZES", "REQUIRES", "LOCATED_AT", "DELIVERS", "MANAGES"',
        "special_rule": "Map operational workflows, logistics, and production."
    },
    "technology": {
        "role": "CTO / Tech Analyst",
        "nodes": '"Technology", "Software", "Hardware", "Protocol", "Innovation"',
        "edges": '"RUNS_ON", "UPGRADES", "DEPRECATES", "ENABLES", "UTILIZES"',
        "special_rule": "Focus on IT infrastructure, digital transformation, and software."
    },
    "strategy": {
        "role": "Chief Strategy Officer",
        "nodes": '"Strategy", "Vision", "Competitor", "Market", "Merger"',
        "edges": '"COMPETES_WITH", "TARGETS", "ACQUIRED", "PLANS_TO", "PARTNERS_WITH"',
        "special_rule": "Extract forward-looking statements, M&A, and competitive landscape."
    },
    "legal": {
        "role": "Legal Counsel",
        "nodes": '"Contract", "Law", "Clause", "Party", "Litigation"',
        "edges": '"SIGNS", "SUES", "ENFORCES", "PROHIBITS", "AGREES_TO"',
        "special_rule": "Focus on contractual obligations, laws, and legal disputes."
    },
    "products": {
        "role": "Product Manager",
        "nodes": '"Product", "Service", "Feature", "Pricing", "CustomerSegment"',
        "edges": '"COSTS", "INCLUDES", "LAUNCHED_IN", "SERVES", "REPLACES"',
        "special_rule": "Catalog products, services, and their attributes."
    },
    "educational": {
        "role": "Academic / Economist",
        "nodes": '"Concept", "Theory", "Indicator", "EconomicEvent", "Definition"',
        "edges": '"CAUSES", "DEFINED_AS", "CORRELATED_WITH", "EXPLAINS", "CONTRADICTS"',
        "special_rule": "Extract definitions, causal relationships, and theoretical concepts."
    },
    "sustainability": {
        "role": "ESG Analyst",
        "nodes": '"ESG_Factor", "Emission", "Initiative", "Standard", "Impact"',
        "edges": '"REDUCES", "IMPACTS", "ALIGNED_WITH", "TARGETS_NET_ZERO", "POLLUTES"',
        "special_rule": "Focus on Environmental, Social, and Governance (ESG) topics."
    },
    "generic": {
        "role": "General Knowledge Analyst",
        "nodes": '"Entity", "Person", "Organization", "Location", "Concept", "Date"',
        "edges": '"RELATED_TO", "INVOLVES", "LOCATED_IN", "HAS_ATTRIBUTE", "MENTIONS"',
        "special_rule": "Extract general entities. Be broad and inclusive."
    }
}


# =========================
# PROMPT TEMPLATES (robusti)
# =========================
CLASSIFIER_PROMPT = """
You are a document classifier for a financial knowledge base.

Return ONLY one token from this exact list:
FINANCIAL, CLIENT, RISK, OPERATIONS, TECHNOLOGY, STRATEGY, LEGAL, PRODUCTS, EDUCATIONAL, SUSTAINABILITY, GENERIC

Rules:
- Output must be exactly one of the tokens above.
- No punctuation, no quotes, no explanations.
- Filename is weak signal, text is primary signal.
- If uncertain or mixed topics, return GENERIC.
"""

MASTER_PROMPT_TEMPLATE = """
ROLE: {role}
TASK: From the TEXT, extract a small Knowledge Graph.

OUTPUT: Return ONLY valid JSON (no markdown) with this schema:
{{
  "nodes": [
    {{"id": "normalized_id", "label": "surface form", "type": "NodeType", "evidence": "short substring"}}
  ],
  "edges": [
    {{"source": "normalized_id", "target": "normalized_id", "relation": "REL_TYPE", "evidence": "short substring"}}
  ]
}}

--- CONTEXT ---
SOURCE FILENAME: "{filename}"
PROCESSING DATE: {today}
TEMPORAL HINT: {temporal_hint}

--- RIGID ONTOLOGY (English only) ---
NODES allowed: {nodes} plus "KeyTerm"
RELATIONS allowed: {edges} plus "RELATED_TO"

--- STRICT RULES ---
1) {special_rule}
2) Use ONLY information explicitly stated in TEXT. Do NOT invent.
3) Every node and edge MUST include "evidence" copied from TEXT (<=120 chars). If you can't cite evidence, omit it.
4) Max nodes: 12. Max edges: 18. No duplicates.
5) "id" MUST be normalized: lowercase, trim, collapse spaces, remove quotes.
6) Keep "label" in original language when possible.

Return JSON only.
"""

VISION_PROMPT = """
You are a financial document analyst.
Extract ONLY what is clearly visible in the image. Do NOT guess.

Return ONLY valid JSON with this schema:
{
  "kind": "chart|table|formula|diagram|photo|other",
  "key_points": ["..."],   // max 6
  "numbers": [{"label":"", "value":"", "unit":"", "currency":"", "period":""}],  // only if clearly readable
  "unreadable_parts": ["..."]  // optional
}

If the image is decorative or unreadable, use kind="other" and empty arrays.
"""


# =========================
# UTILS
# =========================
def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()

def normalize_ws(text: str) -> str:
    # Normalizzazione semplice ma utile: whitespace + line breaks
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def normalize_entity_id(s: str) -> str:
    s = (s or "").strip().lower()
    s = s.replace('"', "").replace("'", "")
    s = re.sub(r"\s+", " ", s)
    return s[:180]  # evita id troppo lunghi

def safe_json_extract(raw: str) -> Optional[dict]:
    """Tenta parsing robusto JSON anche se il modello sputa testo extra."""
    if not raw:
        return None
    s = raw.strip()

    # Rimuovi fence se presenti
    if "```json" in s:
        s = s.split("```json", 1)[1].split("```", 1)[0].strip()
    elif "```" in s:
        s = s.split("```", 1)[1].split("```", 1)[0].strip()

    # Tenta parse diretto
    try:
        return json.loads(s)
    except Exception:
        pass

    # Fallback: estrai primo oggetto JSON {...}
    m = re.search(r"\{.*\}", s, re.DOTALL)
    if m:
        try:
            return json.loads(m.group(0))
        except Exception:
            return None
    return None

def build_dynamic_prompt(doc_type: str, filename: str) -> str:
    ontology = ONTOLOGIES.get(doc_type, ONTOLOGIES["generic"])
    match = re.search(r"(19|20)\d{2}", filename)
    temporal_hint = (
        f"The filename contains '{match.group(0)}'. Assume unstated dates refer to this year."
        if match else "No specific year in filename."
    )
    return MASTER_PROMPT_TEMPLATE.format(
        role=ontology["role"],
        filename=filename,
        today=datetime.datetime.now().strftime("%Y-%m-%d"),
        temporal_hint=temporal_hint,
        nodes=ontology["nodes"],
        edges=ontology["edges"],
        special_rule=ontology["special_rule"]
    )


# =========================
# KG CANDIDATE HEURISTICS
# =========================
DOC_KEYWORDS = {
    "financial": [
        "revenue", "sales", "ebitda", "ebit", "profit", "loss", "cash", "debt",
        "balance sheet", "p&l", "income statement", "net", "gross", "margin",
        "yield", "return", "volatility", "sharpe", "alpha", "beta", "benchmark",
        "aum", "nav", "ter", "fees", "fx", "eur", "usd", "bps"
    ],
    "risk": [
        "risk", "compliance", "aml", "kyc", "audit", "control", "policy", "regulation",
        "violation", "mitigation", "incident", "operational risk", "credit risk"
    ],
    "legal": [
        "shall", "hereby", "agreement", "clause", "party", "liability", "warranty",
        "governing law", "jurisdiction", "termination"
    ],
    "technology": [
        "api", "database", "server", "cloud", "protocol", "architecture", "security",
        "encryption", "deployment", "microservice"
    ],
    "strategy": [
        "strategy", "vision", "roadmap", "merger", "acquisition", "competitor",
        "market share", "growth", "expansion"
    ],
}

NUMERIC_SIGNALS_RE = re.compile(r"(\d{1,3}([.,]\d{3})+|\d+)([.,]\d+)?\s*(%|bps|bp|€|\$|eur|usd|bn|mld|mln|million|billion)?", re.IGNORECASE)

def is_candidate_for_kg(text: str, doc_type: str) -> bool:
    t = (text or "").lower()
    if len(t) < 250:
        # chunk piccoli: richiedi almeno segnali numerici forti
        return bool(NUMERIC_SIGNALS_RE.search(t)) and any(k in t for k in DOC_KEYWORDS.get(doc_type, []))

    if NUMERIC_SIGNALS_RE.search(t):
        return True

    kws = DOC_KEYWORDS.get(doc_type, [])
    if any(k in t for k in kws):
        return True

    # fallback generico: se contiene molte entità stile "Titolo: ..." o ":" frequenti
    if t.count(":") >= 6:
        return True

    return False


# =========================
# SERVIZI: INIT
# =========================
print("🔌 Connessione servizi...")

# Embeddings
try:
    embedder = SentenceTransformer(MODEL_EMBEDDING_ID, device=EMBED_DEVICE)
except Exception:
    print("⚠️ CUDA non disponibile o errore init embedder: fallback su CPU.")
    embedder = SentenceTransformer(MODEL_EMBEDDING_ID, device="cpu")

# Qdrant
qdrant_client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)

# Neo4j
neo4j_driver = GraphDatabase.driver(NEO4J_URI, auth=NEO4J_AUTH)

# LLM client (LM Studio)
llm_client = OpenAI(base_url=LM_STUDIO_URL, api_key=LM_STUDIO_API_KEY)

# Postgres pool
pg_pool = SimpleConnectionPool(PG_POOL_MINCONN, PG_POOL_MAXCONN, dsn=PG_DSN)


# =========================
# POSTGRES HELPERS (pool)
# =========================
def pg_get_conn():
    return pg_pool.getconn()

def pg_put_conn(conn):
    pg_pool.putconn(conn)

def pg_start_log(filename: str, source_type: str) -> int:
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO ingestion_logs (source_name, source_type, status, ingestion_ts) "
                "VALUES (%s, %s, 'PROCESSING', NOW()) RETURNING log_id",
                (filename, source_type)
            )
            log_id = cur.fetchone()[0]
        conn.commit()
        return log_id
    except Exception as e:
        conn.rollback()
        raise e
    finally:
        pg_put_conn(conn)

def pg_close_log(log_id: int, status: str, total_chunks: int, processing_ms: int, error_msg: str = None):
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE ingestion_logs SET status = %s, total_chunks = %s, processing_time_ms = %s, error_message = %s "
                "WHERE log_id = %s",
                (status, total_chunks, processing_ms, error_msg, log_id)
            )
        conn.commit()
    except Exception:
        conn.rollback()
    finally:
        pg_put_conn(conn)

def pg_get_image_by_hash(image_hash: str) -> Optional[Tuple[int, str]]:
    """Ritorna (image_id, description_ai) se esiste."""
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT image_id, description_ai FROM ingestion_images WHERE image_hash = %s LIMIT 1",
                (image_hash,)
            )
            row = cur.fetchone()
            return row if row else None
    except Exception:
        return None
    finally:
        pg_put_conn(conn)

def pg_save_image(log_id: int, image_bytes: bytes, mime_type: str, description: str) -> int:
    """Inserisce solo se non esiste già per hash."""
    img_hash = sha256_hex(image_bytes)

    cached = pg_get_image_by_hash(img_hash)
    if cached:
        return cached[0]

    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO ingestion_images (log_id, image_data, image_hash, mime_type, description_ai, ingestion_ts) "
                "VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING image_id",
                (log_id, psycopg2.Binary(image_bytes), img_hash, mime_type, description)
            )
            image_id = cur.fetchone()[0]
        conn.commit()
        return image_id
    except Exception:
        conn.rollback()
        return -1
    finally:
        pg_put_conn(conn)

def flush_postgres_chunks_batch(batch_data: List[Tuple]):
    if not batch_data:
        return
    conn = pg_get_conn()
    try:
        with conn.cursor() as cur:
            execute_values(
                cur,
                """
                INSERT INTO document_chunks
                    (log_id, chunk_index, toon_type, content_raw, content_semantic, metadata_json)
                VALUES %s
                """,
                batch_data
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        print(f"   ⚠️ Postgres Batch Error: {e}")
    finally:
        pg_put_conn(conn)


# =========================
# NEO4J BATCH (UNWIND)
# =========================
NEO4J_QUERY = """
// 1) Chunk node
MERGE (c:Chunk {id: $chunk_id})
SET c.doc_type = $doc_type,
    c.filename = $filename,
    c.log_id = $log_id,
    c.chunk_index = $chunk_index,
    c.toon_type = $toon_type,
    c.page_no = $page_no,
    c.ingested_at = datetime()

// 2) Entities + mention relationship
WITH c
UNWIND $nodes AS n
MERGE (e:Entity {id: n.id})
ON CREATE SET e.label = n.label, e.type = n.type
ON MATCH  SET e.label = coalesce(e.label, n.label),
             e.type  = coalesce(e.type,  n.type)
MERGE (e)-[:MENTIONED_IN]->(c)

// 3) Relations with provenance
WITH c
UNWIND $edges AS r
MATCH (a:Entity {id: r.source})
MATCH (b:Entity {id: r.target})
MERGE (a)-[rel:REL {type: r.relation, chunk_id: $chunk_id}]->(b)
SET rel.evidence = r.evidence,
    rel.filename = $filename,
    rel.log_id = $log_id
"""

def flush_neo4j_chunk_graph(
    chunk_id: str,
    log_id: int,
    doc_type: str,
    filename: str,
    chunk_index: int,
    toon_type: str,
    page_no: Optional[int],
    graph: Dict
):
    if not graph:
        return

    nodes = graph.get("nodes") or []
    edges = graph.get("edges") or []

    # Sanitize & normalize nodes/edges per sicurezza
    clean_nodes = []
    for n in nodes:
        if not isinstance(n, dict):
            continue
        label = (n.get("label") or n.get("id") or "").strip()
        ntype = (n.get("type") or "KeyTerm").strip()
        evidence = (n.get("evidence") or "").strip()

        nid = normalize_entity_id(n.get("id") or label)
        if not nid or not evidence:
            continue

        clean_nodes.append({
            "id": nid,
            "label": label[:300],
            "type": ntype[:60],
        })

    # Dedup nodes by id
    seen = set()
    dedup_nodes = []
    for n in clean_nodes:
        if n["id"] in seen:
            continue
        seen.add(n["id"])
        dedup_nodes.append(n)

    clean_edges = []
    for e in edges:
        if not isinstance(e, dict):
            continue
        src = normalize_entity_id(e.get("source") or "")
        tgt = normalize_entity_id(e.get("target") or "")
        rel = (e.get("relation") or "RELATED_TO").upper().strip().replace(" ", "_")
        rel = re.sub(r"[^A-Z0-9_]", "", rel) or "RELATED_TO"
        evidence = (e.get("evidence") or "").strip()

        if not src or not tgt or not evidence:
            continue
        # Evita self-loop inutili
        if src == tgt:
            continue

        clean_edges.append({
            "source": src,
            "target": tgt,
            "relation": rel[:80],
            "evidence": evidence[:200]
        })

    # Dedup edges
    seen_e = set()
    dedup_edges = []
    for e in clean_edges:
        key = (e["source"], e["target"], e["relation"])
        if key in seen_e:
            continue
        seen_e.add(key)
        dedup_edges.append(e)

    if not dedup_nodes and not dedup_edges:
        return

    with neo4j_driver.session() as session:
        try:
            session.run(
                NEO4J_QUERY,
                chunk_id=chunk_id,
                log_id=log_id,
                doc_type=doc_type,
                filename=filename,
                chunk_index=chunk_index,
                toon_type=toon_type,
                page_no=page_no,
                nodes=dedup_nodes,
                edges=dedup_edges
            )
        except Exception as e:
            # Non bloccare ingestion per un errore grafo
            print(f"\n   ⚠️ Neo4j write error (chunk {chunk_index}): {e}")


# =========================
# LLM CALLS
# =========================
def classify_document(preview: str, filename: str) -> str:
    if not preview.strip():
        return "generic"

    try:
        resp = llm_client.chat.completions.create(
            model=MODEL_LLM_ID,
            messages=[
                {"role": "system", "content": CLASSIFIER_PROMPT},
                {"role": "user", "content": f"FILENAME:\n{filename}\n\nTEXT PREVIEW:\n{preview[:3000]}"}
            ],
            temperature=0.0,
            max_tokens=5
        )
        token = (resp.choices[0].message.content or "").strip().upper()
        mapping = {
            "FINANCIAL": "financial",
            "CLIENT": "client",
            "RISK": "risk",
            "OPERATIONS": "operations",
            "TECHNOLOGY": "technology",
            "STRATEGY": "strategy",
            "LEGAL": "legal",
            "PRODUCTS": "products",
            "EDUCATIONAL": "educational",
            "SUSTAINABILITY": "sustainability",
            "GENERIC": "generic"
        }
        return mapping.get(token, "generic")
    except Exception as e:
        print(f"⚠️ Errore classificazione AI: {e}")
        return "generic"

def analyze_image_with_vision(image_bytes: bytes) -> str:
    try:
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        resp = llm_client.chat.completions.create(
            model=MODEL_VISION_ID,
            messages=[
                {"role": "system", "content": VISION_PROMPT},
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Analyze the image and return JSON only."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}}
                    ]
                }
            ],
            temperature=0.1,
            max_tokens=500
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception:
        return ""

def extract_graph_data(text: str, doc_type: str, filename: str) -> Optional[Dict]:
    system_prompt = build_dynamic_prompt(doc_type, filename)
    try:
        resp = llm_client.chat.completions.create(
            model=MODEL_LLM_ID,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"TEXT:\n{text}"}
            ],
            temperature=0.0,
            max_tokens=700
        )
        content = (resp.choices[0].message.content or "").strip()
        data = safe_json_extract(content)
        return data
    except Exception:
        return None


# =========================
# PREVIEW EXTRACTION (per classifier)
# =========================
def get_preview_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            doc = fitz.open(file_path)
            out = []
            for i in range(min(PDF_MAX_PAGES_PREVIEW, len(doc))):
                out.append(doc[i].get_text("text"))
            doc.close()
            return "\n".join(out)[:3000]
        elif ext in [".docx", ".doc"]:
            d = Document(file_path)
            return "\n".join(p.text for p in d.paragraphs[:30])[:3000]
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(file_path)
            out = []
            for i, slide in enumerate(prs.slides):
                if i >= 3:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        out.append(shape.text)
            return "\n".join(out)[:3000]
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(3000)
        return ""
    except Exception as e:
        print(f"⚠️ Errore lettura preview: {e}")
        return ""


# =========================
# PARSERS + CHUNKING
# =========================
def split_text_with_overlap(text: str, max_chars: int, overlap: int) -> List[str]:
    text = text.strip()
    if len(text) <= max_chars:
        return [text] if text else []
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        chunk = text[start:end]
        chunks.append(chunk.strip())
        if end == len(text):
            break
        start = max(0, end - overlap)
    return [c for c in chunks if c]

def split_paragraphs(text: str, max_chars: int, overlap: int) -> List[str]:
    """
    Split semantico: accumula paragrafi fino a max_chars, poi crea chunk con overlap (semplice).
    """
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks = []
    buf = ""
    for p in paras:
        if len(buf) + len(p) + 2 <= max_chars:
            buf = (buf + "\n\n" + p).strip() if buf else p
        else:
            if buf:
                chunks.append(buf)
            buf = p

    if buf:
        chunks.append(buf)

    # Applica overlap a livello testo (semplice)
    final = []
    for c in chunks:
        final.extend(split_text_with_overlap(c, max_chars, overlap))
    return [c for c in final if c]

def remove_repeated_lines(pages_text: List[str], min_len: int = 10, max_len: int = 80, freq_ratio: float = 0.6) -> List[str]:
    """
    Rimuove righe che compaiono su molte pagine (tipici header/footer).
    """
    if len(pages_text) < 4:
        return pages_text

    line_counts: Dict[str, int] = {}
    for t in pages_text:
        lines = set([ln.strip() for ln in t.splitlines() if min_len <= len(ln.strip()) <= max_len])
        for ln in lines:
            line_counts[ln] = line_counts.get(ln, 0) + 1

    threshold = int(len(pages_text) * freq_ratio)
    repeated = {ln for ln, c in line_counts.items() if c >= threshold}

    if not repeated:
        return pages_text

    cleaned = []
    for t in pages_text:
        lines = [ln for ln in t.splitlines() if ln.strip() not in repeated]
        cleaned.append("\n".join(lines))
    return cleaned

def extract_pdf_chunks(file_path: str, log_id: int) -> List[Dict]:
    """
    Estrae chunk da PDF:
    - 1 livello: testo per pagina (con pulizia righe ripetute)
    - 2 livello: split semantico per paragrafi + overlap
    - Vision: solo immagini grandi + limitate; caching su hash; opzionale solo se testo pagina scarso
    Restituisce lista di dict: {text_raw, text_sem, page_no, toon_type, images_meta...}
    """
    doc = fitz.open(file_path)
    pages_text = [doc[i].get_text("text") for i in range(len(doc))]
    pages_text = [normalize_ws(t) for t in pages_text]
    pages_text = remove_repeated_lines(pages_text)

    out_chunks: List[Dict] = []
    print(f"   🔍 Analisi PDF ({len(doc)} pagine)...")

    for page_no in range(len(doc)):
        page = doc[page_no]
        page_text = pages_text[page_no]
        page_text_len = len(page_text)

        # Vision gating
        allow_vision = True
        if PDF_VISION_ONLY_IF_TEXT_SCARSO:
            allow_vision = page_text_len < PDF_TEXT_SCARSO_THRESHOLD

        # Estrai immagini (limitate)
        image_desc_blocks = []
        if allow_vision:
            images = page.get_images(full=True) or []
            used = 0
            for img_info in images:
                if used >= PDF_VISION_MAX_IMAGES_PER_PAGE:
                    break
                xref = img_info[0]
                try:
                    base = doc.extract_image(xref)
                    img_bytes = base.get("image", b"")
                    ext = base.get("ext", "jpg")

                    if not img_bytes or len(img_bytes) < PDF_IMAGE_MIN_BYTES:
                        continue

                    # dimensioni immagine (se disponibili)
                    # img_info: (xref, smask, width, height, bpc, colorspace, alt, name, filter)
                    width = img_info[2] if len(img_info) > 3 else 0
                    height = img_info[3] if len(img_info) > 3 else 0
                    if width and height and (width < PDF_IMAGE_MIN_DIM or height < PDF_IMAGE_MIN_DIM):
                        continue

                    # Cache: se hash già presente, riusa descrizione
                    img_hash = sha256_hex(img_bytes)
                    cached = pg_get_image_by_hash(img_hash)
                    if cached:
                        img_id, desc_ai = cached[0], cached[1] or ""
                    else:
                        desc_ai = analyze_image_with_vision(img_bytes)
                        # prova a parsare JSON, se non valido -> fallback testo
                        parsed = safe_json_extract(desc_ai)
                        if parsed and parsed.get("kind") == "other" and not parsed.get("key_points") and not parsed.get("numbers"):
                            # non rilevante
                            continue
                        img_id = pg_save_image(log_id, img_bytes, f"image/{ext}", desc_ai)

                    if img_id > 0:
                        image_desc_blocks.append(f"[[IMAGE_DESC_AI (ID: {img_id}, PAGE: {page_no+1}): {desc_ai}]]")
                        used += 1
                except Exception:
                    continue

        # Costruisci testo per pagina con eventuali blocchi immagini
        combined_raw = page_text
        if image_desc_blocks:
            combined_raw += "\n\n" + "\n\n".join(image_desc_blocks)

        combined_sem = normalize_ws(combined_raw)

        # Split in chunk semantici
        page_chunks = split_paragraphs(combined_sem, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)
        for ch in page_chunks:
            toon_type = "image_description" if "[[IMAGE_DESC_AI" in ch else "text"
            out_chunks.append({
                "text_raw": ch,
                "text_sem": ch,   # qui puoi differenziare se vuoi: raw vs semantic; al momento ch è già "sem"
                "page_no": page_no + 1,
                "toon_type": toon_type
            })

    doc.close()
    return out_chunks

def extract_docx_chunks(file_path: str) -> List[Dict]:
    d = Document(file_path)
    paras = [p.text.strip() for p in d.paragraphs if p.text and p.text.strip()]
    text = "\n\n".join(paras)
    text = normalize_ws(text)
    chunks = split_paragraphs(text, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)
    return [{"text_raw": c, "text_sem": c, "page_no": None, "toon_type": "text"} for c in chunks]

def extract_pptx_chunks(file_path: str) -> List[Dict]:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                buf.append(shape.text.strip())
        s = normalize_ws("\n".join(buf))
        if s:
            slides.append((i + 1, s))

    out = []
    for slide_no, s in slides:
        chunks = split_text_with_overlap(s, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)
        for c in chunks:
            out.append({"text_raw": c, "text_sem": c, "page_no": slide_no, "toon_type": "text"})
    return out

def extract_txt_chunks(file_path: str) -> List[Dict]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = normalize_ws(f.read())
    chunks = split_paragraphs(text, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS)
    return [{"text_raw": c, "text_sem": c, "page_no": None, "toon_type": "text"} for c in chunks]

def extract_file_chunks(file_path: str, log_id: int) -> List[Dict]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_chunks(file_path, log_id)
    elif ext in [".docx", ".doc"]:
        return extract_docx_chunks(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_pptx_chunks(file_path)
    elif ext == ".txt":
        return extract_txt_chunks(file_path)
    return []


# =========================
# QDRANT INIT
# =========================
def ensure_qdrant_collection():
    if not qdrant_client.collection_exists(QDRANT_COLLECTION):
        print(f"🧱 Creo collection Qdrant: {QDRANT_COLLECTION}")
        qdrant_client.create_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=models.VectorParams(size=1024, distance=models.Distance.COSINE)
        )


# =========================
# ARCHIVE
# =========================
def archive_file(file_path: str, category: str):
    try:
        dest = os.path.join(PROCESSED_DIR, category)
        os.makedirs(dest, exist_ok=True)
        fname = os.path.basename(file_path)
        name, ext = os.path.splitext(fname)
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        shutil.move(file_path, os.path.join(dest, f"{name}_{ts}{ext}"))
        print(f"   📦 Archived to {category}")
    except Exception as e:
        print(f"   ❌ Archive Error: {e}")


# =========================
# MAIN ENGINE
# =========================
def process_single_file(file_path: str, doc_type: str):
    filename = os.path.basename(file_path)
    start_time = time.time()
    print(f"   ⚙️ Engine Start: {filename}")

    # Start log
    try:
        log_id = pg_start_log(filename, doc_type)
    except Exception as e:
        print(f"❌ DB Start Error: {e}")
        return

    try:
        # Estrai chunks (page-aware per PDF)
        chunks = extract_file_chunks(file_path, log_id)
        chunks = [c for c in chunks if c.get("text_sem") and len(c["text_sem"]) >= MIN_CHUNK_LEN]

        if not chunks:
            elapsed = int((time.time() - start_time) * 1000)
            pg_close_log(log_id, "SKIPPED_EMPTY", 0, elapsed)
            return

        print(f"   🚀 Estratti {len(chunks)} chunk (doc_type={doc_type}).")
        print(f"   🧠 Embeddings batch_size={EMBED_BATCH_SIZE}, flush_size={DB_FLUSH_SIZE}, max_KG={MAX_KG_CHUNKS_PER_DOC}")

        # Buffers
        qdrant_points: List[models.PointStruct] = []
        pg_rows: List[Tuple] = []
        neo4j_jobs: List[Tuple] = []  # (chunk_uuid, chunk_index, toon_type, page_no, graph)

        # KG budget
        kg_done = 0

        # Process in embedding batches
        total_saved = 0
        chunk_index_global = 0

        for batch_start in range(0, len(chunks), EMBED_BATCH_SIZE):
            batch = chunks[batch_start: batch_start + EMBED_BATCH_SIZE]
            texts_sem = [b["text_sem"] for b in batch]

            # 1) Embeddings in batch
            vectors = embedder.encode(texts_sem, batch_size=EMBED_BATCH_SIZE, normalize_embeddings=True)

            # 2) Per chunk: prepare Qdrant + PG + (optional) KG
            for local_i, ch in enumerate(batch):
                text_raw = ch["text_raw"]
                text_sem = ch["text_sem"]
                toon_type = ch.get("toon_type") or "text"
                page_no = ch.get("page_no")

                chunk_uuid = str(uuid.uuid4())
                vec = vectors[local_i].tolist()

                # Qdrant
                payload = {
                    "filename": filename,
                    "type": doc_type,
                    "chunk_index": chunk_index_global,
                    "toon_type": toon_type,
                    "page_no": page_no,
                    "log_id": log_id
                }
                qdrant_points.append(models.PointStruct(id=chunk_uuid, vector=vec, payload=payload))

                # Postgres
                meta = {
                    "qdrant_uuid": chunk_uuid,
                    "original_path": file_path,
                    "filename": filename,
                    "doc_type": doc_type,
                    "chunk_index": chunk_index_global,
                    "page_no": page_no,
                    "toon_type": toon_type
                }
                pg_rows.append((log_id, chunk_index_global, toon_type, text_raw, text_sem, Json(meta)))

                # KG extraction: solo se candidato + budget disponibile
                if kg_done < MAX_KG_CHUNKS_PER_DOC and is_candidate_for_kg(text_sem, doc_type):
                    graph = extract_graph_data(text_sem, doc_type, filename)
                    if graph:
                        neo4j_jobs.append((chunk_uuid, chunk_index_global, toon_type, page_no, graph))
                        kg_done += 1

                total_saved += 1
                chunk_index_global += 1
                print(".", end="", flush=True)

                # Flush if needed
                if len(qdrant_points) >= DB_FLUSH_SIZE:
                    # Qdrant
                    qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_points)
                    # Postgres
                    flush_postgres_chunks_batch(pg_rows)
                    # Neo4j (per chunk, ma con UNWIND interno; in alternativa: raggruppa ancora di più)
                    for job in neo4j_jobs:
                        cid, cidx, ttype, pno, g = job
                        flush_neo4j_chunk_graph(cid, log_id, doc_type, filename, cidx, ttype, pno, g)

                    qdrant_points.clear()
                    pg_rows.clear()
                    neo4j_jobs.clear()
                    print(" [Flush] ", end="", flush=True)

        # Final flush
        if qdrant_points:
            qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=qdrant_points)
            flush_postgres_chunks_batch(pg_rows)
            for job in neo4j_jobs:
                cid, cidx, ttype, pno, g = job
                flush_neo4j_chunk_graph(cid, log_id, doc_type, filename, cidx, ttype, pno, g)
            print(" [Flush] ", end="", flush=True)

        elapsed = int((time.time() - start_time) * 1000)
        pg_close_log(log_id, "COMPLETED", total_saved, elapsed)
        print(f"\n✅ Completed: {filename} | chunks={total_saved} | KG_calls={kg_done} | ms={elapsed}")

        archive_file(file_path, doc_type)

    except Exception as e:
        elapsed = int((time.time() - start_time) * 1000)
        print(f"\n❌ Critical Error: {e}")
        pg_close_log(log_id, "FAILED", 0, elapsed, str(e))


def main():
    os.makedirs(INBOX_DIR, exist_ok=True)
    os.makedirs(PROCESSED_DIR, exist_ok=True)

    ensure_qdrant_collection()

    files = glob.glob(os.path.join(INBOX_DIR, "*"))
    files = [f for f in files if os.path.isfile(f) and not f.endswith(".tmp")]

    if not files:
        print(f"📭 INBOX vuota ({INBOX_DIR})")
        return

    print(f"📬 Trovati {len(files)} file in INBOX.")

    for f in files:
        fname = os.path.basename(f)
        print(f"\n🕵️ Classificazione: {fname}...")

        preview = get_preview_text(f)
        doc_type = classify_document(preview, fname)

        print(f"   🏷️ Categoria rilevata: {doc_type.upper()}")
        process_single_file(f, doc_type)

    print("\n✅ Ingestion Finished.")

    try:
        neo4j_driver.close()
    except Exception:
        pass

    try:
        pg_pool.closeall()
    except Exception:
        pass


if __name__ == "__main__":
    main()
